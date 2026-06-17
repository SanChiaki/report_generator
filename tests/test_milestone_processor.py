from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Inches

from report_generator.generator import generate_report
from report_generator.post_processing import PostProcessingRegistry
from report_generator.pptx.document import PptxDocument


def _template_bytes() -> bytes:
    from io import BytesIO

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    anchor = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.7),
        Inches(1.0),
        Inches(8.0),
        Inches(1.2),
    )
    anchor.name = "milestone.delivery"
    anchor.text = "milestone anchor"
    anchor.fill.solid()
    anchor.fill.fore_color.rgb = RGBColor(0xFF, 0xF2, 0xCC)
    output = BytesIO()
    prs.save(output)
    return output.getvalue()


def _theme_template_bytes() -> bytes:
    from io import BytesIO

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    anchor = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.7),
        Inches(1.0),
        Inches(8.0),
        Inches(1.2),
    )
    anchor.name = "milestone.delivery"
    anchor.fill.solid()
    anchor.fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_4
    anchor.fill.fore_color.brightness = 0.8
    output = BytesIO()
    prs.save(output)
    return output.getvalue()


def _preview_parts_template_bytes() -> bytes:
    from io import BytesIO

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    anchor = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.7),
        Inches(1.0),
        Inches(8.0),
        Inches(1.2),
    )
    anchor.name = "milestone.delivery"
    anchor.fill.solid()
    anchor.fill.fore_color.rgb = RGBColor(0xFF, 0xF2, 0xCC)
    preview_node = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(1.0),
        Inches(1.45),
        Inches(0.2),
        Inches(0.2),
    )
    preview_node.name = "milestone.delivery.preview.node_1"
    output = BytesIO()
    prs.save(output)
    return output.getvalue()


def test_milestone_draws_dynamic_nodes_from_items():
    mapping = {
        "template_id": "milestone-test",
        "component_list": [
            {
                "location": "milestone.delivery",
                "semantic_description": "项目里程碑",
                "type": "Milestone",
                "data_source": {"name": "milestones"},
            }
        ],
    }
    payload = {
        "milestones": {
            "items": [
                {"label": "准备", "date": "04-20", "status": "done"},
                {"label": "安装", "date": "05-18", "status": "done"},
                {"label": "调测", "date": "06-20", "status": "active"},
                {"label": "验收", "date": "07-05", "status": "pending"},
                {"label": "上线", "date": "07-20", "status": "pending"},
            ]
        }
    }

    output = generate_report(_template_bytes(), mapping, payload, PostProcessingRegistry())
    doc = PptxDocument.open(output)
    index = doc.shape_index()

    assert index["milestone.delivery"].shape.text == ""
    assert index["milestone.delivery"].shape.fill.fore_color.rgb == RGBColor(0xFF, 0xF2, 0xCC)
    nodes = [ref for name, ref in index.items() if name.startswith("milestone.delivery.item_") and name.endswith(".node")]
    dates = [ref for name, ref in index.items() if name.startswith("milestone.delivery.item_") and name.endswith(".date")]
    labels = [ref for name, ref in index.items() if name.startswith("milestone.delivery.item_") and name.endswith(".label")]
    assert len(nodes) == 5
    assert [ref.shape.text for ref in dates] == ["04-20", "05-18", "06-20", "07-05", "07-20"]
    assert [ref.shape.text for ref in labels] == ["准备", "安装", "调测", "验收", "上线"]
    assert nodes[1].shape.left > nodes[0].shape.left
    assert nodes[4].shape.left > nodes[3].shape.left


def test_milestone_preview_mode_replace_removes_template_preview_shape():
    mapping = {
        "template_id": "milestone-test",
        "component_list": [
            {
                "location": "milestone.delivery",
                "semantic_description": "项目里程碑",
                "type": "Milestone",
                "config": {"preview_mode": "replace"},
                "data_source": {"name": "milestones"},
            }
        ],
    }
    payload = {"milestones": [{"label": "上线", "date": "07-20", "status": "done"}]}

    output = generate_report(_template_bytes(), mapping, payload, PostProcessingRegistry())
    doc = PptxDocument.open(output)
    index = doc.shape_index()

    assert "milestone.delivery" not in index
    assert index["milestone.delivery.item_1.label"].shape.text == "上线"


def test_milestone_preview_mode_replace_recreates_anchor_background():
    mapping = {
        "template_id": "milestone-test",
        "component_list": [
            {
                "location": "milestone.delivery",
                "semantic_description": "项目里程碑",
                "type": "Milestone",
                "config": {"preview_mode": "replace"},
                "data_source": {"name": "milestones"},
            }
        ],
    }
    payload = {"milestones": [{"label": "上线", "date": "07-20", "status": "done"}]}

    output = generate_report(_template_bytes(), mapping, payload, PostProcessingRegistry())
    doc = PptxDocument.open(output)
    background = doc.shape_index()["milestone.delivery.background"].shape

    assert background.fill.fore_color.rgb == RGBColor(0xFF, 0xF2, 0xCC)
    assert background.left == Inches(0.7)
    assert background.top == Inches(1.0)
    assert background.width == Inches(8.0)
    assert background.height == Inches(1.2)


def test_milestone_preview_mode_replace_preserves_theme_anchor_background():
    mapping = {
        "template_id": "milestone-test",
        "component_list": [
            {
                "location": "milestone.delivery",
                "semantic_description": "项目里程碑",
                "type": "Milestone",
                "config": {"preview_mode": "replace"},
                "data_source": {"name": "milestones"},
            }
        ],
    }
    payload = {"milestones": [{"label": "上线", "date": "07-20", "status": "done"}]}

    output = generate_report(_theme_template_bytes(), mapping, payload, PostProcessingRegistry())
    doc = PptxDocument.open(output)
    background = doc.shape_index()["milestone.delivery.background"].shape

    assert background.fill.fore_color.theme_color == MSO_THEME_COLOR.ACCENT_4
    assert background.fill.fore_color.brightness == 0.8


def test_milestone_preview_mode_replace_removes_named_preview_parts():
    mapping = {
        "template_id": "milestone-test",
        "component_list": [
            {
                "location": "milestone.delivery",
                "semantic_description": "项目里程碑",
                "type": "Milestone",
                "config": {"preview_mode": "replace"},
                "data_source": {"name": "milestones"},
            }
        ],
    }
    payload = {"milestones": [{"label": "上线", "date": "07-20", "status": "done"}]}

    output = generate_report(_preview_parts_template_bytes(), mapping, payload, PostProcessingRegistry())
    doc = PptxDocument.open(output)
    index = doc.shape_index()

    assert "milestone.delivery.preview.node_1" not in index
    assert "milestone.delivery.background" in index
    assert index["milestone.delivery.item_1.label"].shape.text == "上线"


def test_milestone_uses_readable_default_text_sizes_and_label_font():
    mapping = {
        "template_id": "milestone-test",
        "component_list": [
            {
                "location": "milestone.delivery",
                "semantic_description": "项目里程碑",
                "type": "Milestone",
                "data_source": {"name": "milestones"},
            }
        ],
    }
    payload = {"milestones": [{"label": "准备", "date": "04-20", "status": "done"}]}

    output = generate_report(_template_bytes(), mapping, payload, PostProcessingRegistry())
    doc = PptxDocument.open(output)
    index = doc.shape_index()
    date_run = index["milestone.delivery.item_1.date"].shape.text_frame.paragraphs[0].runs[0]
    label_run = index["milestone.delivery.item_1.label"].shape.text_frame.paragraphs[0].runs[0]

    assert date_run.font.size.pt == 14
    assert _typefaces(date_run) == {
        "latin": "Microsoft YaHei",
        "ea": "Microsoft YaHei",
        "cs": "Microsoft YaHei",
    }
    assert label_run.font.size.pt == 14
    assert _typefaces(label_run) == {
        "latin": "Microsoft YaHei",
        "ea": "Microsoft YaHei",
        "cs": "Microsoft YaHei",
    }
    assert label_run.font.bold is True


def test_milestone_accepts_plain_list_and_centers_single_node():
    mapping = {
        "template_id": "milestone-test",
        "component_list": [
            {
                "location": "milestone.delivery",
                "semantic_description": "项目里程碑",
                "type": "Milestone",
                "data_source": {"name": "milestones"},
            }
        ],
    }
    payload = {"milestones": [{"label": "上线", "date": "07-20", "status": "done"}]}

    output = generate_report(_template_bytes(), mapping, payload, PostProcessingRegistry())
    doc = PptxDocument.open(output)
    index = doc.shape_index()

    node = index["milestone.delivery.item_1.node"].shape
    label = index["milestone.delivery.item_1.label"].shape
    assert label.text == "上线"
    assert abs((node.left + node.width / 2) - Inches(4.7)) < Inches(0.05)


def test_milestone_treats_numeric_dimensions_as_inches():
    mapping = {
        "template_id": "milestone-test",
        "component_list": [
            {
                "location": "milestone.delivery",
                "semantic_description": "项目里程碑",
                "type": "Milestone",
                "config": {"node_size": 0.25, "date_height": 0.3, "label_height": 0.3},
                "data_source": {"name": "milestones"},
            }
        ],
    }
    payload = {"milestones": [{"label": "上线", "date": "07-20", "status": "done"}]}

    output = generate_report(_template_bytes(), mapping, payload, PostProcessingRegistry())
    doc = PptxDocument.open(output)
    index = doc.shape_index()

    assert index["milestone.delivery.item_1.node"].shape.width == Inches(0.25)
    assert index["milestone.delivery.item_1.date"].shape.height == Inches(0.3)
    assert index["milestone.delivery.item_1.label"].shape.height == Inches(0.3)


def test_milestone_draws_nodes_without_outline_by_default():
    mapping = {
        "template_id": "milestone-test",
        "component_list": [
            {
                "location": "milestone.delivery",
                "semantic_description": "项目里程碑",
                "type": "Milestone",
                "data_source": {"name": "milestones"},
            }
        ],
    }
    payload = {"milestones": [{"label": "上线", "date": "07-20", "status": "done"}]}

    output = generate_report(_template_bytes(), mapping, payload, PostProcessingRegistry())
    doc = PptxDocument.open(output)
    node = doc.shape_index()["milestone.delivery.item_1.node"].shape

    assert node.line.width.pt == 0


def test_milestone_supports_configured_node_outline_width():
    mapping = {
        "template_id": "milestone-test",
        "component_list": [
            {
                "location": "milestone.delivery",
                "semantic_description": "项目里程碑",
                "type": "Milestone",
                "config": {"node_outline_width": 0.75},
                "data_source": {"name": "milestones"},
            }
        ],
    }
    payload = {"milestones": [{"label": "准备", "date": "04-20", "status": "done"}]}

    output = generate_report(_template_bytes(), mapping, payload, PostProcessingRegistry())
    doc = PptxDocument.open(output)
    node = doc.shape_index()["milestone.delivery.item_1.node"].shape

    assert node.line.width.pt == 0.75
    assert node.line.fill.fore_color.rgb == node.fill.fore_color.rgb


def test_milestone_draws_hollow_nodes_as_two_filled_circles():
    mapping = {
        "template_id": "milestone-test",
        "component_list": [
            {
                "location": "milestone.delivery",
                "semantic_description": "项目里程碑",
                "type": "Milestone",
                "data_source": {"name": "milestones"},
            }
        ],
    }
    payload = {"milestones": [{"label": "验收", "date": "07-05", "status": "pending"}]}

    output = generate_report(_template_bytes(), mapping, payload, PostProcessingRegistry())
    doc = PptxDocument.open(output)
    index = doc.shape_index()
    outer = index["milestone.delivery.item_1.node"].shape
    inner = index["milestone.delivery.item_1.node.inner"].shape

    assert outer.fill.fore_color.rgb == RGBColor(0x9A, 0xA6, 0xB2)
    assert outer.line.width.pt == 0
    assert inner.fill.fore_color.rgb == RGBColor(0xFF, 0xFF, 0xFF)
    assert inner.line.width.pt == 0
    assert inner.width < outer.width
    assert inner.left + inner.width / 2 == outer.left + outer.width / 2


def test_milestone_supports_horizontal_padding_date_color_and_centered_text():
    mapping = {
        "template_id": "milestone-test",
        "component_list": [
            {
                "location": "milestone.delivery",
                "semantic_description": "项目里程碑",
                "type": "Milestone",
                "config": {"horizontal_padding": 1, "date_color": "C00000"},
                "data_source": {"name": "milestones"},
            }
        ],
    }
    payload = {
        "milestones": [
            {"label": "准备", "date": "04-20", "status": "done"},
            {"label": "验收", "date": "07-05", "status": "pending"},
        ]
    }

    output = generate_report(_template_bytes(), mapping, payload, PostProcessingRegistry())
    doc = PptxDocument.open(output)
    index = doc.shape_index()
    first_node = index["milestone.delivery.item_1.node"].shape
    second_node = index["milestone.delivery.item_2.node"].shape
    first_date = index["milestone.delivery.item_1.date"].shape
    first_label = index["milestone.delivery.item_1.label"].shape
    second_date = index["milestone.delivery.item_2.date"].shape
    second_label = index["milestone.delivery.item_2.label"].shape
    date_run = first_date.text_frame.paragraphs[0].runs[0]
    first_node_center = first_node.left + first_node.width / 2
    second_node_center = second_node.left + second_node.width / 2

    assert abs((first_node.left + first_node.width / 2) - Inches(1.7)) < Inches(0.05)
    assert abs((second_node.left + second_node.width / 2) - Inches(7.7)) < Inches(0.05)
    assert first_date.left + first_date.width / 2 == first_node_center
    assert first_label.left + first_label.width / 2 == first_node_center
    assert second_date.left + second_date.width / 2 == second_node_center
    assert second_label.left + second_label.width / 2 == second_node_center
    assert date_run.font.color.rgb == RGBColor(0xC0, 0x00, 0x00)


def test_milestone_vertically_middle_aligns_date_and_label_text():
    mapping = {
        "template_id": "milestone-test",
        "component_list": [
            {
                "location": "milestone.delivery",
                "semantic_description": "项目里程碑",
                "type": "Milestone",
                "config": {"date_height": 0.32, "label_height": 0.36},
                "data_source": {"name": "milestones"},
            }
        ],
    }
    payload = {"milestones": [{"label": "验收", "date": "2026-07-20", "status": "done"}]}

    output = generate_report(_template_bytes(), mapping, payload, PostProcessingRegistry())
    doc = PptxDocument.open(output)
    index = doc.shape_index()
    date = index["milestone.delivery.item_1.date"].shape
    label = index["milestone.delivery.item_1.label"].shape

    assert date.text_frame.vertical_anchor == MSO_ANCHOR.MIDDLE
    assert label.text_frame.vertical_anchor == MSO_ANCHOR.MIDDLE


def test_milestone_supports_custom_date_width_for_full_dates():
    mapping = {
        "template_id": "milestone-test",
        "component_list": [
            {
                "location": "milestone.delivery",
                "semantic_description": "项目里程碑",
                "type": "Milestone",
                "config": {"date_width": 1.1},
                "data_source": {"name": "milestones"},
            }
        ],
    }
    payload = {"milestones": [{"label": "上线", "date": "2026-07-20", "status": "done"}]}

    output = generate_report(_template_bytes(), mapping, payload, PostProcessingRegistry())
    doc = PptxDocument.open(output)
    index = doc.shape_index()
    node = index["milestone.delivery.item_1.node"].shape
    date = index["milestone.delivery.item_1.date"].shape

    assert date.width == Inches(1.1)
    assert date.left + date.width / 2 == node.left + node.width / 2


def test_milestone_keeps_date_and_label_equally_spaced_from_axis():
    mapping = {
        "template_id": "milestone-test",
        "component_list": [
            {
                "location": "milestone.delivery",
                "semantic_description": "项目里程碑",
                "type": "Milestone",
                "config": {
                    "node_size": 0.2,
                    "date_height": 0.24,
                    "label_height": 0.3,
                    "text_axis_gap": 0.14,
                },
                "data_source": {"name": "milestones"},
            }
        ],
    }
    payload = {"milestones": [{"label": "上线", "date": "07-20", "status": "done"}]}

    output = generate_report(_template_bytes(), mapping, payload, PostProcessingRegistry())
    doc = PptxDocument.open(output)
    index = doc.shape_index()
    date = index["milestone.delivery.item_1.date"].shape
    label = index["milestone.delivery.item_1.label"].shape
    node = index["milestone.delivery.item_1.node"].shape
    axis_y = node.top + node.height / 2

    assert axis_y - (date.top + date.height) == Inches(0.14)
    assert label.top - axis_y == Inches(0.14)


def _typefaces(run):
    r_pr = run._r.rPr
    values = {}
    for child in r_pr:
        tag = child.tag
        if tag.endswith(("}latin", "}ea", "}cs")):
            values[tag.split("}", 1)[1]] = child.get("typeface")
    return values
