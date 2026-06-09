import pytest
from pptx.dml.color import RGBColor

from report_generator.components.table import apply_table
from report_generator.errors import ErrorCode, ReportGenerationError
from report_generator.models import ComponentMapping
from report_generator.pptx.document import PptxDocument


def table_component(**config):
    return ComponentMapping.model_validate(
        {
            "location": "table.top_risks",
            "semantic_description": "TOP问题与风险",
            "type": "Table",
            "config": config,
        }
    )


def test_apply_table_rebuilds_rows_and_columns(simple_template_bytes):
    doc = PptxDocument.open(simple_template_bytes)
    shape = doc.shape_index()["table.top_risks"].shape
    data = [
        {"风险类型": "延期风险", "风险描述": "设备到货存在延期风险"},
        {"风险类型": "质量风险", "风险描述": "测试通过率偏低"},
    ]

    new_shape = apply_table(
        doc,
        shape,
        table_component(order=["风险类型", "风险描述"], min_font_size=8),
        data,
    )

    table = new_shape.table
    assert len(table.rows) == 3
    assert len(table.columns) == 2
    assert table.cell(0, 0).text == "风险类型"
    assert table.cell(2, 1).text == "测试通过率偏低"


def test_apply_table_defaults_to_dynamic_resize_while_copying_template_styles(simple_template_bytes):
    doc = PptxDocument.open(simple_template_bytes)
    shape = doc.shape_index()["table.top_risks"].shape
    shape.table.cell(0, 0).fill.solid()
    shape.table.cell(0, 0).fill.fore_color.rgb = RGBColor(0x12, 0x34, 0x56)
    shape.table.cell(1, 0).fill.solid()
    shape.table.cell(1, 0).fill.fore_color.rgb = RGBColor(0xAB, 0xCD, 0xEF)
    shape.table.cell(0, 0).text_frame.paragraphs[0].runs[0].font.size = None
    shape.table.cell(1, 0).text_frame.paragraphs[0].runs[0].font.bold = True
    header_fill = shape.table.cell(0, 0).fill.fore_color.rgb
    body_fill = shape.table.cell(1, 0).fill.fore_color.rgb
    body_bold = shape.table.cell(1, 0).text_frame.paragraphs[0].runs[0].font.bold
    data = [
        {"风险类型": "延期风险", "风险描述": "设备到货存在延期风险"},
        {"风险类型": "质量风险", "风险描述": "测试通过率偏低"},
    ]

    new_shape = apply_table(
        doc,
        shape,
        table_component(order=["风险类型", "风险描述"], min_font_size=8),
        data,
    )

    assert len(new_shape.table.rows) == 3
    assert new_shape.table.cell(0, 0).fill.fore_color.rgb == header_fill
    assert new_shape.table.cell(1, 0).fill.fore_color.rgb == body_fill
    assert new_shape.table.cell(2, 0).fill.fore_color.rgb == body_fill
    assert new_shape.table.cell(1, 0).text_frame.paragraphs[0].runs[0].font.bold == body_bold
    assert new_shape.table.cell(2, 1).text == "测试通过率偏低"


def test_apply_table_appends_rows_in_place_when_columns_fit(simple_template_bytes):
    doc = PptxDocument.open(simple_template_bytes)
    shape = doc.shape_index()["table.top_risks"].shape
    original_shape_id = shape.shape_id
    original_height = shape.height
    template_body_height = shape.table.rows[1].height
    shape.table.cell(1, 0).fill.solid()
    shape.table.cell(1, 0).fill.fore_color.rgb = RGBColor(0xAB, 0xCD, 0xEF)
    shape.table.cell(1, 0).text_frame.paragraphs[0].runs[0].font.bold = True
    body_fill = shape.table.cell(1, 0).fill.fore_color.rgb
    body_bold = shape.table.cell(1, 0).text_frame.paragraphs[0].runs[0].font.bold
    data = [
        {"风险类型": "延期风险", "风险描述": "设备到货存在延期风险"},
        {"风险类型": "质量风险", "风险描述": "测试通过率偏低"},
        {"风险类型": "资源风险", "风险描述": "关键人员排期冲突"},
    ]

    updated_shape = apply_table(
        doc,
        shape,
        table_component(order=["风险类型", "风险描述"], max_rows=3),
        data,
    )

    assert updated_shape.shape_id == original_shape_id
    assert len(updated_shape.table.rows) == 4
    assert updated_shape.height == original_height + (2 * template_body_height)
    assert updated_shape.table.cell(3, 0).text == "资源风险"
    assert updated_shape.table.cell(3, 0).fill.fore_color.rgb == body_fill
    assert updated_shape.table.cell(3, 0).text_frame.paragraphs[0].runs[0].font.bold == body_bold


def test_apply_table_rejects_too_many_rows(simple_template_bytes):
    doc = PptxDocument.open(simple_template_bytes)
    shape = doc.shape_index()["table.top_risks"].shape
    data = [{"风险类型": str(index), "风险描述": "x"} for index in range(3)]

    with pytest.raises(ReportGenerationError) as exc:
        apply_table(
            doc,
            shape,
            table_component(order=["风险类型", "风险描述"], max_rows=2),
            data,
        )

    assert exc.value.error_code == ErrorCode.TABLE_OVERFLOW


def test_apply_table_accepts_columns_rows_object(simple_template_bytes):
    doc = PptxDocument.open(simple_template_bytes)
    shape = doc.shape_index()["table.top_risks"].shape
    value = {
        "columns": [
            {"key": "type", "label": "风险类型"},
            {"key": "desc", "label": "风险描述"},
        ],
        "rows": [{"type": "延期风险", "desc": "设备到货存在延期风险"}],
    }

    new_shape = apply_table(doc, shape, table_component(), value)

    assert new_shape.table.cell(0, 0).text == "风险类型"
    assert new_shape.table.cell(1, 1).text == "设备到货存在延期风险"


def test_apply_table_accepts_cells_matrix_without_header(simple_template_bytes):
    doc = PptxDocument.open(simple_template_bytes)
    shape = doc.shape_index()["table.top_risks"].shape
    shape.table.cell(0, 0).text_frame.paragraphs[0].runs[0].font.bold = True
    shape.table.cell(0, 1).text = ""
    shape.table.cell(1, 1).fill.solid()
    shape.table.cell(1, 1).fill.fore_color.rgb = RGBColor(0xAB, 0xCD, 0xEF)
    body_fill = shape.table.cell(1, 1).fill.fore_color.rgb
    value = {
        "cells": [
            ["项目名称", "智慧园区", "项目编码", "PRJ-2025-001"],
            ["客户名称", "示例客户集团", "项目 PD", "张明"],
        ]
    }

    new_shape = apply_table(doc, shape, table_component(), value)

    assert len(new_shape.table.rows) == 2
    assert len(new_shape.table.columns) == 4
    assert new_shape.table.cell(0, 0).text == "项目名称"
    assert new_shape.table.cell(0, 1).text == "智慧园区"
    assert new_shape.table.cell(0, 1).text_frame.paragraphs[0].runs[0].font.bold is True
    assert new_shape.table.cell(1, 3).text == "张明"
    assert new_shape.table.cell(1, 3).fill.fore_color.rgb == body_fill


def test_apply_table_preserves_existing_table_shape_and_style(simple_template_bytes):
    doc = PptxDocument.open(simple_template_bytes)
    shape = doc.shape_index()["table.top_risks"].shape
    original_shape_id = shape.shape_id
    shape.table.cell(0, 0).fill.solid()
    shape.table.cell(0, 0).fill.fore_color.rgb = RGBColor(0x12, 0x34, 0x56)
    original_header_fill = shape.table.cell(0, 0).fill.fore_color.rgb
    value = {
        "columns": [
            {"key": "type", "label": "风险类型"},
            {"key": "desc", "label": "风险描述"},
        ],
        "rows": [{"type": "延期风险", "desc": "设备到货存在延期风险"}],
    }

    updated_shape = apply_table(doc, shape, table_component(preserve_style=True), value)

    assert updated_shape.shape_id == original_shape_id
    assert updated_shape.table.cell(0, 0).text == "风险类型"
    assert updated_shape.table.cell(1, 1).text == "设备到货存在延期风险"
    assert updated_shape.table.cell(0, 0).fill.fore_color.rgb == original_header_fill


def test_apply_table_preserve_style_accepts_cells_matrix(simple_template_bytes):
    doc = PptxDocument.open(simple_template_bytes)
    shape = doc.shape_index()["table.top_risks"].shape
    original_shape_id = shape.shape_id
    value = {"cells": [["项目名称", "智慧园区"], ["客户名称", "示例客户集团"]]}

    updated_shape = apply_table(doc, shape, table_component(preserve_style=True), value)

    assert updated_shape.shape_id == original_shape_id
    assert updated_shape.table.cell(0, 0).text == "项目名称"
    assert updated_shape.table.cell(0, 1).text == "智慧园区"
    assert updated_shape.table.cell(1, 0).text == "客户名称"
    assert updated_shape.table.cell(1, 1).text == "示例客户集团"


def test_apply_table_replaces_cell_placeholders_without_rebuilding(simple_template_bytes):
    doc = PptxDocument.open(simple_template_bytes)
    shape = doc.shape_index()["table.top_risks"].shape
    original_shape_id = shape.shape_id
    shape.table.cell(0, 1).text = "{{ 项目名称 }}"
    shape.table.cell(0, 0).text_frame.paragraphs[0].runs[0].font.bold = True

    updated_shape = apply_table(
        doc,
        shape,
        table_component(mode="placeholders"),
        {"项目名称": "智慧园区"},
    )

    assert updated_shape.shape_id == original_shape_id
    assert len(updated_shape.table.rows) == 2
    assert len(updated_shape.table.columns) == 2
    assert updated_shape.table.cell(0, 0).text == "风险类型"
    assert updated_shape.table.cell(0, 1).text == "智慧园区"
    assert updated_shape.table.cell(0, 1).text_frame.paragraphs[0].runs[0].font.bold is True


def test_apply_table_preserve_style_rejects_when_template_has_insufficient_rows(simple_template_bytes):
    doc = PptxDocument.open(simple_template_bytes)
    shape = doc.shape_index()["table.top_risks"].shape
    data = [
        {"风险类型": "延期风险", "风险描述": "设备到货存在延期风险"},
        {"风险类型": "质量风险", "风险描述": "测试通过率偏低"},
    ]

    with pytest.raises(ReportGenerationError) as exc:
        apply_table(
            doc,
            shape,
            table_component(order=["风险类型", "风险描述"], preserve_style=True),
            data,
        )

    assert exc.value.error_code == ErrorCode.TABLE_OVERFLOW


def test_apply_table_rejects_non_object_rows(simple_template_bytes):
    doc = PptxDocument.open(simple_template_bytes)
    shape = doc.shape_index()["table.top_risks"].shape

    with pytest.raises(ReportGenerationError) as exc:
        apply_table(doc, shape, table_component(order=["风险类型"]), ["not-a-row"])

    assert exc.value.error_code == ErrorCode.DATA_SOURCE_INVALID


def test_apply_table_rejects_non_object_columns(simple_template_bytes):
    doc = PptxDocument.open(simple_template_bytes)
    shape = doc.shape_index()["table.top_risks"].shape
    value = {"columns": ["bad"], "rows": [{"bad": "value"}]}

    with pytest.raises(ReportGenerationError) as exc:
        apply_table(doc, shape, table_component(), value)

    assert exc.value.error_code == ErrorCode.DATA_SOURCE_INVALID
