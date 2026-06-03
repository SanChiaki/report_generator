from __future__ import annotations

from io import BytesIO

import pytest
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt


def _save(prs: Presentation) -> bytes:
    output = BytesIO()
    prs.save(output)
    return output.getvalue()


@pytest.fixture
def simple_template_bytes() -> bytes:
    prs = Presentation()
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(6), Inches(0.6))
    title.name = "text.report_title"
    title.text_frame.text = "旧标题"
    title.text_frame.paragraphs[0].runs[0].font.size = Pt(24)

    summary = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(4), Inches(0.7))
    summary.name = "text.summary"
    summary.text_frame.text = "旧摘要"
    summary.text_frame.paragraphs[0].runs[0].font.size = Pt(18)

    table_shape = slide.shapes.add_table(2, 2, Inches(0.5), Inches(2), Inches(5), Inches(1.2))
    table_shape.name = "table.top_risks"
    table = table_shape.table
    table.cell(0, 0).text = "风险类型"
    table.cell(0, 1).text = "风险描述"
    table.cell(1, 0).text = "旧类型"
    table.cell(1, 1).text = "旧描述"

    chart_data = CategoryChartData()
    chart_data.categories = ["Q1", "Q2"]
    chart_data.add_series("收入", (10, 20))
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(6),
        Inches(2),
        Inches(3),
        Inches(2),
        chart_data,
    )
    chart.name = "chart.revenue_trend"

    badge = slide.shapes.add_shape(1, Inches(6), Inches(0.5), Inches(1.2), Inches(0.5))
    badge.name = "shape.status_badge"
    badge.text = "旧状态"

    image_box = slide.shapes.add_textbox(Inches(6), Inches(1.1), Inches(1.2), Inches(0.7))
    image_box.name = "image.company_logo"
    image_box.text = "logo"

    return _save(prs)


@pytest.fixture
def duplicate_name_template_bytes() -> bytes:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    first = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(2), Inches(0.5))
    second = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(2), Inches(0.5))
    first.name = "text.duplicate"
    second.name = "text.duplicate"
    return _save(prs)
