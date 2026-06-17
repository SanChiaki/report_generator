from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from report_generator.generator import generate_report
from report_generator.post_processing import PostProcessingRegistry
from report_generator.pptx.document import PptxDocument


def _template_bytes() -> bytes:
    from io import BytesIO

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    anchor = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5),
        Inches(0.8),
        Inches(8.5),
        Inches(1.0),
    )
    anchor.name = "top_issues.cards"
    anchor.text = "TOP issues anchor"
    output = BytesIO()
    prs.save(output)
    return output.getvalue()


def _preview_parts_template_bytes() -> bytes:
    from io import BytesIO

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    anchor = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5),
        Inches(0.8),
        Inches(8.5),
        Inches(1.0),
    )
    anchor.name = "top_issues.cards"
    anchor.fill.solid()
    anchor.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    preview_card = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5),
        Inches(0.8),
        Inches(8.5),
        Inches(1.0),
    )
    preview_card.name = "top_issues.cards.preview.card"
    preview_card.fill.solid()
    preview_card.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    output = BytesIO()
    prs.save(output)
    return output.getvalue()


def test_top_issues_draws_dynamic_vertical_cards_with_severity_styles():
    mapping = {
        "template_id": "top-issues-test",
        "component_list": [
            {
                "location": "top_issues.cards",
                "semantic_description": "TOP 问题与风险",
                "type": "TopIssues",
                "config": {"card_height": 0.6, "card_gap": 0.1},
                "data_source": {"name": "issues"},
            }
        ],
    }
    payload = {
        "issues": [
            {
                "severity": "紧急",
                "created_at": "2026-06-10",
                "description": "验收窗口未确认",
                "action": "协调客户确认备选时间",
                "owner": "张明",
                "status": "跟踪中",
                "due_date": "2026-06-14",
            },
            {
                "severity": "重要",
                "description": "兼容性补测未完成",
                "action": "补充自动化用例",
                "owner": "王琳",
                "status": "处理中",
                "due_date": "2026-06-16",
            },
            {
                "severity": "一般",
                "description": "指标口径待确认",
                "action": "周会对齐",
                "owner": "赵强",
                "status": "待确认",
                "due_date": "2026-06-18",
            },
            {
                "severity": "紧急",
                "description": "额外问题",
                "action": "继续跟踪",
                "owner": "张明",
                "status": "新增",
                "due_date": "2026-06-20",
            },
        ]
    }

    output = generate_report(_template_bytes(), mapping, payload, PostProcessingRegistry())
    doc = PptxDocument.open(output)
    index = doc.shape_index()

    assert "top_issues.cards" not in index
    cards = [ref for name, ref in index.items() if name.startswith("top_issues.cards.item_") and name.endswith(".card")]
    strips = [ref for name, ref in index.items() if name.startswith("top_issues.cards.item_") and name.endswith(".strip")]
    severities = [
        ref for name, ref in index.items() if name.startswith("top_issues.cards.item_") and name.endswith(".severity")
    ]
    assert len(cards) == 4
    assert len(strips) == 4
    assert [ref.shape.text for ref in severities] == ["紧急", "重要", "一般", "紧急"]
    assert strips[0].shape.fill.fore_color.rgb == RGBColor(0xD6, 0x45, 0x45)
    assert strips[1].shape.fill.fore_color.rgb == RGBColor(0xFF, 0x8A, 0x00)
    assert strips[2].shape.fill.fore_color.rgb == RGBColor(0xF5, 0xC4, 0x00)
    assert cards[1].shape.top > cards[0].shape.top
    assert cards[3].shape.top > cards[2].shape.top


def test_top_issues_preview_mode_replace_removes_template_preview_shape():
    mapping = {
        "template_id": "top-issues-test",
        "component_list": [
            {
                "location": "top_issues.cards",
                "semantic_description": "TOP 问题与风险",
                "type": "TopIssues",
                "config": {"preview_mode": "replace"},
                "data_source": {"name": "issues"},
            }
        ],
    }
    payload = {"issues": [{"severity": "紧急", "description": "验收窗口未确认"}]}

    output = generate_report(_template_bytes(), mapping, payload, PostProcessingRegistry())
    doc = PptxDocument.open(output)
    index = doc.shape_index()

    assert "top_issues.cards" not in index
    assert index["top_issues.cards.item_1.description"].shape.text == "问题描述：验收窗口未确认"


def test_top_issues_preview_mode_replace_removes_named_preview_parts():
    mapping = {
        "template_id": "top-issues-test",
        "component_list": [
            {
                "location": "top_issues.cards",
                "semantic_description": "TOP 问题与风险",
                "type": "TopIssues",
                "config": {"preview_mode": "replace"},
                "data_source": {"name": "issues"},
            }
        ],
    }
    payload = {"issues": [{"severity": "紧急", "description": "验收窗口未确认"}]}

    output = generate_report(_preview_parts_template_bytes(), mapping, payload, PostProcessingRegistry())
    doc = PptxDocument.open(output)
    index = doc.shape_index()

    assert "top_issues.cards.preview.card" not in index
    assert "top_issues.cards.item_1.card" in index


def test_top_issues_accepts_items_object_and_custom_templates():
    mapping = {
        "template_id": "top-issues-test",
        "component_list": [
            {
                "location": "top_issues.cards",
                "semantic_description": "TOP 问题与风险",
                "type": "TopIssues",
                "config": {
                    "description_template": "问题：{{ title }}",
                    "action_template": "措施：{{ action }}",
                    "meta_template": "负责人：{{ owner }}",
                },
                "data_source": {"name": "issues"},
            }
        ],
    }
    payload = {
        "issues": {
            "items": [
                {
                    "severity": "重要",
                    "title": "客户验收排期待确认",
                    "action": "锁定评审窗口",
                    "owner": "张明",
                }
            ]
        }
    }

    output = generate_report(_template_bytes(), mapping, payload, PostProcessingRegistry())
    doc = PptxDocument.open(output)
    texts = [ref.shape.text for ref in doc.shape_index().values() if getattr(ref.shape, "has_text_frame", False)]

    assert "问题：客户验收排期待确认" in texts
    assert "措施：锁定评审窗口" in texts
    assert "负责人：张明" in texts


def test_top_issues_uses_template_like_default_font_sizes():
    mapping = {
        "template_id": "top-issues-test",
        "component_list": [
            {
                "location": "top_issues.cards",
                "semantic_description": "TOP 问题与风险",
                "type": "TopIssues",
                "data_source": {"name": "issues"},
            }
        ],
    }
    payload = {
        "issues": [
            {
                "severity": "紧急",
                "created_at": "2026-06-10",
                "description": "验收窗口未确认",
                "action": "协调客户确认备选时间",
                "owner": "张明",
                "status": "跟踪中",
                "due_date": "2026-06-14",
            }
        ]
    }

    output = generate_report(_template_bytes(), mapping, payload, PostProcessingRegistry())
    doc = PptxDocument.open(output)
    index = doc.shape_index()
    severity_run = index["top_issues.cards.item_1.severity"].shape.text_frame.paragraphs[0].runs[0]
    description_run = index["top_issues.cards.item_1.description"].shape.text_frame.paragraphs[0].runs[0]
    action_run = index["top_issues.cards.item_1.action"].shape.text_frame.paragraphs[0].runs[0]
    meta_run = index["top_issues.cards.item_1.meta"].shape.text_frame.paragraphs[0].runs[0]

    assert severity_run.font.size.pt == 18
    assert description_run.font.size.pt == 13
    assert action_run.font.size.pt == 13
    assert meta_run.font.size.pt == 12


def test_top_issues_uses_microsoft_yahei_for_all_generated_text():
    mapping = {
        "template_id": "top-issues-test",
        "component_list": [
            {
                "location": "top_issues.cards",
                "semantic_description": "TOP 问题与风险",
                "type": "TopIssues",
                "data_source": {"name": "issues"},
            }
        ],
    }
    payload = {
        "issues": [
            {
                "severity": "紧急",
                "description": "验收窗口未确认",
                "action": "协调客户确认备选时间",
                "owner": "Owner A",
                "status": "tracking",
                "due_date": "2026-06-14",
            }
        ]
    }

    output = generate_report(_template_bytes(), mapping, payload, PostProcessingRegistry())
    doc = PptxDocument.open(output)
    index = doc.shape_index()
    runs = [
        index["top_issues.cards.item_1.severity"].shape.text_frame.paragraphs[0].runs[0],
        *index["top_issues.cards.item_1.description"].shape.text_frame.paragraphs[0].runs,
        *index["top_issues.cards.item_1.action"].shape.text_frame.paragraphs[0].runs,
        index["top_issues.cards.item_1.meta"].shape.text_frame.paragraphs[0].runs[0],
    ]

    for run in runs:
        assert _typefaces(run) == {
            "latin": "Microsoft YaHei",
            "ea": "Microsoft YaHei",
            "cs": "Microsoft YaHei",
        }


def test_top_issues_styles_description_and_action_labels_separately_from_values():
    mapping = {
        "template_id": "top-issues-test",
        "component_list": [
            {
                "location": "top_issues.cards",
                "semantic_description": "TOP 问题与风险",
                "type": "TopIssues",
                "config": {
                    "description_color": "0052CC",
                    "action_color": "0052CC",
                },
                "data_source": {"name": "issues"},
            }
        ],
    }
    payload = {
        "issues": [
            {
                "severity": "紧急",
                "description": "验收窗口未确认",
                "action": "协调客户确认备选时间",
            }
        ]
    }

    output = generate_report(_template_bytes(), mapping, payload, PostProcessingRegistry())
    doc = PptxDocument.open(output)
    index = doc.shape_index()
    description_runs = index["top_issues.cards.item_1.description"].shape.text_frame.paragraphs[0].runs
    action_runs = index["top_issues.cards.item_1.action"].shape.text_frame.paragraphs[0].runs

    assert [run.text for run in description_runs] == ["问题描述：", "验收窗口未确认"]
    assert description_runs[0].font.bold is True
    assert description_runs[0].font.color.rgb == RGBColor(0x00, 0x52, 0xCC)
    assert description_runs[1].font.bold is False
    assert description_runs[1].font.color.rgb == RGBColor(0x33, 0x33, 0x33)
    assert [run.text for run in action_runs] == ["解决措施与进展：", "协调客户确认备选时间"]
    assert action_runs[0].font.bold is True
    assert action_runs[0].font.color.rgb == RGBColor(0x00, 0x52, 0xCC)
    assert action_runs[1].font.bold is False
    assert action_runs[1].font.color.rgb == RGBColor(0x33, 0x33, 0x33)


def test_top_issues_treats_numeric_dimensions_as_inches():
    mapping = {
        "template_id": "top-issues-test",
        "component_list": [
            {
                "location": "top_issues.cards",
                "semantic_description": "TOP 问题与风险",
                "type": "TopIssues",
                "config": {"card_height": 1, "card_gap": 0.25},
                "data_source": {"name": "issues"},
            }
        ],
    }
    payload = {
        "issues": [
            {"severity": "紧急", "description": "第一个问题"},
            {"severity": "重要", "description": "第二个问题"},
        ]
    }

    output = generate_report(_template_bytes(), mapping, payload, PostProcessingRegistry())
    doc = PptxDocument.open(output)
    index = doc.shape_index()
    first = index["top_issues.cards.item_1.card"].shape
    second = index["top_issues.cards.item_2.card"].shape

    assert first.height == Inches(1)
    assert second.top - first.top == Inches(1.25)


def _typefaces(run):
    r_pr = run._r.rPr
    values = {}
    for child in r_pr:
        tag = child.tag
        if tag.endswith(("}latin", "}ea", "}cs")):
            values[tag.split("}", 1)[1]] = child.get("typeface")
    return values
