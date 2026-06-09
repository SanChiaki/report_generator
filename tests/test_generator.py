import pytest
import time
from threading import Lock
from pptx import Presentation
from pptx.util import Inches

from report_generator.errors import ErrorCode, ReportGenerationError
from report_generator.generator import generate_report
from report_generator.post_processing import PostProcessingRegistry
from report_generator.pptx.document import PptxDocument


class FakeComponentProcessor:
    def __init__(self, result):
        self.result = result
        self.calls = []

    def process(self, component, value):
        self.calls.append((component, value))
        return self.result


class BlockingComponentProcessor:
    def __init__(self):
        self._lock = Lock()
        self.active = 0
        self.max_active = 0
        self.starts = []

    def process(self, component, value):
        with self._lock:
            self.active += 1
            self.max_active = max(self.max_active, self.active)
            self.starts.append(component.location)
        time.sleep(0.05)
        with self._lock:
            self.active -= 1
        return f"{component.location}:{value}"


def test_generate_report_populates_text_and_table(simple_template_bytes):
    mapping = {
        "template_id": "project-monthly-report-ppt-v1",
        "component_list": [
            {
                "location": "text.report_title",
                "semantic_description": "报告标题",
                "type": "Text",
                "data_source": {
                    "name": "api_data",
                    "template": "{{ 全量数据.项目概览.项目名称 }}-报告",
                },
            },
            {
                "location": "table.top_risks",
                "semantic_description": "TOP问题与风险",
                "type": "Table",
                "config": {"order": ["风险类型", "风险描述"]},
                "data_source": {
                    "name": "risks",
                },
            },
        ],
    }
    payload = {
        "api_data": {"全量数据": {"项目概览": {"项目名称": "智慧园区"}}},
        "risks": [{"风险类型": "延期风险", "风险描述": "设备到货存在延期风险"}],
    }

    output = generate_report(simple_template_bytes, mapping, payload, PostProcessingRegistry())
    doc = PptxDocument.open(output)
    index = doc.shape_index()

    assert index["text.report_title"].shape.text_frame.text == "智慧园区-报告"
    assert index["table.top_risks"].shape.table.cell(1, 1).text == "设备到货存在延期风险"


def test_generate_report_uses_component_processor_for_post_processing(simple_template_bytes):
    mapping = {
        "template_id": "project-monthly-report-ppt-v1",
        "component_list": [
            {
                "location": "text.report_title",
                "semantic_description": "报告标题",
                "prompt": "把原始项目数据改写成标题",
                "data_example": "智慧园区项目周报",
                "type": "Text",
                "data_source": {
                    "name": "raw_project",
                    "needs_post_processing": True,
                },
            }
        ],
    }
    payload = {"raw_project": {"name": "智慧园区", "status": "green"}}
    processor = FakeComponentProcessor("智慧园区项目周报")

    output = generate_report(
        simple_template_bytes,
        mapping,
        payload,
        PostProcessingRegistry(),
        processor,
    )
    doc = PptxDocument.open(output)

    assert doc.shape_index()["text.report_title"].shape.text_frame.text == "智慧园区项目周报"
    assert processor.calls[0][1] == {"name": "智慧园区", "status": "green"}


def test_generate_report_processes_post_processing_components_in_parallel(simple_template_bytes):
    mapping = {
        "template_id": "project-monthly-report-ppt-v1",
        "component_list": [
            {
                "location": "text.report_title",
                "semantic_description": "报告标题",
                "type": "Text",
                "data_source": {
                    "name": "title",
                    "needs_post_processing": True,
                },
            },
            {
                "location": "text.summary",
                "semantic_description": "摘要",
                "type": "Text",
                "data_source": {
                    "name": "summary",
                    "needs_post_processing": True,
                },
            },
        ],
    }
    processor = BlockingComponentProcessor()

    output = generate_report(
        simple_template_bytes,
        mapping,
        {"title": "标题原始数据", "summary": "摘要原始数据"},
        PostProcessingRegistry(),
        processor,
        llm_concurrency=2,
    )
    doc = PptxDocument.open(output)

    assert processor.max_active == 2
    assert doc.shape_index()["text.report_title"].shape.text_frame.text == "text.report_title:标题原始数据"
    assert doc.shape_index()["text.summary"].shape.text_frame.text == "text.summary:摘要原始数据"


def test_generate_report_limits_post_processing_concurrency(simple_template_bytes):
    mapping = {
        "template_id": "project-monthly-report-ppt-v1",
        "component_list": [
            {
                "location": "text.report_title",
                "semantic_description": "报告标题",
                "type": "Text",
                "data_source": {
                    "name": "title",
                    "needs_post_processing": True,
                },
            },
            {
                "location": "text.summary",
                "semantic_description": "摘要",
                "type": "Text",
                "data_source": {
                    "name": "summary",
                    "needs_post_processing": True,
                },
            },
        ],
    }
    processor = BlockingComponentProcessor()

    generate_report(
        simple_template_bytes,
        mapping,
        {"title": "标题原始数据", "summary": "摘要原始数据"},
        PostProcessingRegistry(),
        processor,
        llm_concurrency=1,
    )

    assert processor.max_active == 1


def test_generate_report_rejects_missing_component(simple_template_bytes):
    mapping = {
        "template_id": "project-monthly-report-ppt-v1",
        "component_list": [
            {
                "location": "text.missing",
                "semantic_description": "缺失组件",
                "type": "Text",
                "data_source": {"name": "api_data"},
            }
        ],
    }

    with pytest.raises(ReportGenerationError) as exc:
        generate_report(simple_template_bytes, mapping, {"api_data": "value"}, PostProcessingRegistry())

    assert exc.value.error_code == ErrorCode.COMPONENT_NOT_FOUND


def test_generate_report_rejects_missing_component_before_post_processing(simple_template_bytes):
    mapping = {
        "template_id": "project-monthly-report-ppt-v1",
        "component_list": [
            {
                "location": "text.missing",
                "semantic_description": "缺失组件",
                "type": "Text",
                "data_source": {
                    "name": "raw_project",
                    "needs_post_processing": True,
                },
            }
        ],
    }
    processor = FakeComponentProcessor("不应调用")

    with pytest.raises(ReportGenerationError) as exc:
        generate_report(
            simple_template_bytes,
            mapping,
            {"raw_project": {"name": "智慧园区"}},
            PostProcessingRegistry(),
            processor,
            llm_concurrency=2,
        )

    assert exc.value.error_code == ErrorCode.COMPONENT_NOT_FOUND
    assert processor.calls == []


def test_generate_report_ignores_duplicate_unmapped_shape_names():
    prs = Presentation()
    for slide_index in range(2):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        mapped = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(4), Inches(0.5))
        mapped.name = f"text.title_{slide_index}"
        mapped.text = "旧标题"
        unmapped = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(4), Inches(0.5))
        unmapped.name = "TextBox 2"
        unmapped.text = "未映射"

    template = _presentation_bytes(prs)
    mapping = {
        "template_id": "project-monthly-report-ppt-v1",
        "component_list": [
            {
                "location": "text.title_0",
                "semantic_description": "第一页标题",
                "type": "Text",
                "data_source": {"name": "title"},
            }
        ],
    }

    output = generate_report(template, mapping, {"title": "新标题"}, PostProcessingRegistry())
    doc = PptxDocument.open(output)

    assert doc.shape_index(required_names={"text.title_0"})["text.title_0"].shape.text_frame.text == "新标题"


def _presentation_bytes(prs: Presentation) -> bytes:
    from io import BytesIO

    output = BytesIO()
    prs.save(output)
    return output.getvalue()
