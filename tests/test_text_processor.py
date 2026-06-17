import pytest
from pptx.dml.color import RGBColor
from pptx.util import Pt

from report_generator.components.text import apply_text
from report_generator.errors import ErrorCode, ReportGenerationError
from report_generator.models import ComponentMapping
from report_generator.pptx.document import PptxDocument


def text_component(**config):
    return ComponentMapping.model_validate(
        {
            "location": "text.summary",
            "semantic_description": "整体状态",
            "type": "Text",
            "config": config,
        }
    )


def test_apply_text_replaces_content(simple_template_bytes):
    doc = PptxDocument.open(simple_template_bytes)
    shape = doc.shape_index()["text.summary"].shape

    apply_text(shape, text_component(min_font_size=10), "新的摘要")

    assert shape.text_frame.text == "新的摘要"


def test_apply_text_preserves_existing_run_style(simple_template_bytes):
    doc = PptxDocument.open(simple_template_bytes)
    shape = doc.shape_index()["text.summary"].shape
    original_run = shape.text_frame.paragraphs[0].runs[0]
    original_size = original_run.font.size

    apply_text(shape, text_component(preserve_style=True), "保留样式的新摘要")

    updated_run = shape.text_frame.paragraphs[0].runs[0]
    assert shape.text_frame.text == "保留样式的新摘要"
    assert updated_run.font.size == original_size


def test_apply_text_preserve_style_expands_height_when_original_font_cannot_fit(simple_template_bytes):
    doc = PptxDocument.open(simple_template_bytes)
    shape = doc.shape_index()["text.summary"].shape
    original_run = shape.text_frame.paragraphs[0].runs[0]
    original_size = original_run.font.size
    original_width = shape.width
    original_height = shape.height

    apply_text(shape, text_component(preserve_style=True), "超长内容" * 80)

    updated_run = shape.text_frame.paragraphs[0].runs[0]
    assert shape.width == original_width
    assert shape.height > original_height
    assert updated_run.font.size == original_size


def test_apply_text_shrinks_long_content(simple_template_bytes):
    doc = PptxDocument.open(simple_template_bytes)
    shape = doc.shape_index()["text.summary"].shape

    apply_text(shape, text_component(min_font_size=8), "这是一段较长的摘要内容，用来触发字号缩小逻辑。" * 3)

    assert shape.text_frame.paragraphs[0].runs[0].font.size <= Pt(18)
    assert shape.text_frame.paragraphs[0].runs[0].font.size >= Pt(8)


def test_apply_text_expands_height_when_content_cannot_fit_original_box(simple_template_bytes):
    doc = PptxDocument.open(simple_template_bytes)
    shape = doc.shape_index()["text.summary"].shape
    original_width = shape.width
    original_height = shape.height

    apply_text(shape, text_component(min_font_size=16), "超长内容" * 400)

    assert shape.width == original_width
    assert shape.height > original_height
    assert shape.text_frame.paragraphs[0].runs[0].font.size == Pt(16)


def test_apply_text_supports_rich_text_runs_with_explicit_styles(simple_template_bytes):
    doc = PptxDocument.open(simple_template_bytes)
    shape = doc.shape_index()["text.summary"].shape

    apply_text(
        shape,
        text_component(min_font_size=10),
        {
            "rich_text": [
                {
                    "text": "关键结论：",
                    "color": "0052CC",
                    "font_size": 16,
                    "font_name": "Microsoft YaHei",
                    "bold": True,
                },
                {
                    "text": "整体推进中",
                    "color": "333333",
                    "font_size": 14,
                    "font_name": "Microsoft YaHei",
                    "bold": False,
                },
            ]
        },
    )

    paragraph = shape.text_frame.paragraphs[0]
    assert shape.text_frame.text == "关键结论：整体推进中"
    assert [run.text for run in paragraph.runs] == ["关键结论：", "整体推进中"]
    assert paragraph.runs[0].font.color.rgb == RGBColor(0x00, 0x52, 0xCC)
    assert paragraph.runs[0].font.size == Pt(16)
    assert paragraph.runs[0].font.name == "Microsoft YaHei"
    assert paragraph.runs[0].font.bold is True
    assert paragraph.runs[1].font.color.rgb == RGBColor(0x33, 0x33, 0x33)
    assert paragraph.runs[1].font.size == Pt(14)
    assert paragraph.runs[1].font.bold is False


def test_apply_text_accepts_runs_alias_for_rich_text(simple_template_bytes):
    doc = PptxDocument.open(simple_template_bytes)
    shape = doc.shape_index()["text.summary"].shape

    apply_text(
        shape,
        text_component(min_font_size=10),
        {
            "runs": [
                {"text": "A", "color": "0052CC"},
                {"text": "B", "color": "333333"},
            ]
        },
    )

    paragraph = shape.text_frame.paragraphs[0]
    assert [run.text for run in paragraph.runs] == ["A", "B"]
    assert paragraph.runs[0].font.color.rgb == RGBColor(0x00, 0x52, 0xCC)
    assert paragraph.runs[1].font.color.rgb == RGBColor(0x33, 0x33, 0x33)


def test_apply_text_supports_newlines_inside_rich_text(simple_template_bytes):
    doc = PptxDocument.open(simple_template_bytes)
    shape = doc.shape_index()["text.summary"].shape

    apply_text(
        shape,
        text_component(min_font_size=10),
        {
            "rich_text": [
                {"text": "关键结论：", "color": "0052CC", "bold": True},
                {"text": "第一行\n第二行", "color": "333333"},
            ]
        },
    )

    assert shape.text_frame.text == "关键结论：第一行\n第二行"
    assert [run.text for run in shape.text_frame.paragraphs[0].runs] == ["关键结论：", "第一行"]
    assert [run.text for run in shape.text_frame.paragraphs[1].runs] == ["第二行"]
    assert shape.text_frame.paragraphs[0].runs[0].font.color.rgb == RGBColor(0x00, 0x52, 0xCC)
    assert shape.text_frame.paragraphs[1].runs[0].font.color.rgb == RGBColor(0x33, 0x33, 0x33)
