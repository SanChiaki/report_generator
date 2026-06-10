from __future__ import annotations

from collections.abc import Mapping
from copy import deepcopy
from dataclasses import dataclass
from typing import Any

from jinja2 import StrictUndefined
from jinja2.exceptions import TemplateError, TemplateSyntaxError, UndefinedError
from jinja2.sandbox import SandboxedEnvironment
from pptx.dml.color import RGBColor
from pptx.util import Pt

from report_generator.errors import ErrorCode, ReportGenerationError
from report_generator.models import ComponentMapping
from report_generator.pptx.document import PptxDocument


@dataclass(frozen=True)
class NormalizedTable:
    cells: list[list[Any]]
    data_row_count: int
    kind: str

    @property
    def row_count(self) -> int:
        return len(self.cells)

    @property
    def column_count(self) -> int:
        return max((len(row) for row in self.cells), default=0)

    def value_at(self, row_index: int, col_index: int) -> Any:
        row = self.cells[row_index]
        if col_index >= len(row):
            return ""
        return row[col_index]


def apply_table(
    doc: PptxDocument,
    shape: Any,
    component: ComponentMapping,
    value: Any,
) -> Any:
    if not getattr(shape, "has_table", False):
        raise ReportGenerationError(
            ErrorCode.TYPE_MISMATCH,
            f"组件 {component.location} 不是表格组件",
            component,
        )

    if component.config.get("mode") in {"placeholder", "placeholders"}:
        return _apply_table_placeholders(shape, component, value)

    table_data = _normalize_table(value, component)
    max_rows = int(component.config.get("max_rows", 30))
    max_columns = int(component.config.get("max_columns", 10))
    if table_data.data_row_count > max_rows:
        raise ReportGenerationError(
            ErrorCode.TABLE_OVERFLOW,
            f"组件 {component.location} 的数据行数 {table_data.data_row_count} 超过限制 {max_rows}",
            component,
        )
    if table_data.column_count > max_columns:
        raise ReportGenerationError(
            ErrorCode.TABLE_OVERFLOW,
            f"组件 {component.location} 的列数 {table_data.column_count} 超过限制 {max_columns}",
            component,
        )

    if component.config.get("preserve_style"):
        updated_shape = _apply_table_preserving_style(shape, component, table_data)
        _apply_column_value_styles(updated_shape.table, component, table_data)
        return updated_shape

    if table_data.column_count <= len(shape.table.columns):
        updated_shape = _apply_table_with_row_append(shape, component, table_data)
        _apply_column_value_styles(updated_shape.table, component, table_data)
        return updated_shape

    updated_shape = _rebuild_table(doc, shape, component, table_data)
    _apply_column_value_styles(updated_shape.table, component, table_data)
    return updated_shape


def _rebuild_table(
    doc: PptxDocument,
    shape: Any,
    component: ComponentMapping,
    table_data: NormalizedTable,
) -> Any:
    row_count = table_data.row_count
    column_count = table_data.column_count
    if row_count == 0 or column_count == 0:
        raise ReportGenerationError(
            ErrorCode.DATA_SOURCE_INVALID,
            f"组件 {component.location} 的表格数据为空且无法推断列",
            component,
        )

    x, y, cx, cy = shape.left, shape.top, shape.width, shape.height
    source_table = shape.table
    slide = shape.part.slide
    new_shape = slide.shapes.add_table(row_count, column_count, x, y, cx, cy)
    new_shape.name = component.location
    table = new_shape.table
    _copy_table_style(source_table, table, cx, cy)
    doc.remove_shape(shape)

    _write_cells(table, table_data)

    if "font_size" in component.config:
        _apply_font_size(table, int(component.config["font_size"]))

    return new_shape


def _apply_table_with_row_append(
    shape: Any,
    component: ComponentMapping,
    table_data: NormalizedTable,
) -> Any:
    table = shape.table
    row_count = table_data.row_count
    column_count = table_data.column_count
    if row_count == 0 or column_count == 0:
        raise ReportGenerationError(
            ErrorCode.DATA_SOURCE_INVALID,
            f"组件 {component.location} 的表格数据为空且无法推断列",
            component,
        )
    if column_count > len(table.columns):
        raise ReportGenerationError(
            ErrorCode.TABLE_OVERFLOW,
            f"组件 {component.location} 的列数 {column_count} 超过模板预留列数 {len(table.columns)}",
            component,
        )

    while len(table.rows) < row_count:
        _append_table_row(table)

    _write_cells(table, table_data)
    for row_index in range(row_count, len(table.rows)):
        for col_index in range(len(table.columns)):
            _replace_cell_text_preserving_style(table.cell(row_index, col_index), "")
    for row_index in range(row_count):
        for col_index in range(column_count, len(table.columns)):
            _replace_cell_text_preserving_style(table.cell(row_index, col_index), "")

    shape.height = sum(table.rows[row_index].height for row_index in range(len(table.rows)))
    if "font_size" in component.config:
        _apply_font_size(table, int(component.config["font_size"]))

    return shape


def _append_table_row(table: Any) -> None:
    source_row = table.rows[len(table.rows) - 1]
    table._tbl.append(deepcopy(source_row._tr))


def _copy_table_style(source_table: Any, target_table: Any, width: int, height: int) -> None:
    _copy_dimensions(source_table, target_table, width, height)
    source_row_count = len(source_table.rows)
    source_col_count = len(source_table.columns)
    for row_index in range(len(target_table.rows)):
        source_row_index = _source_index(row_index, source_row_count, repeat_after_first=True)
        for col_index in range(len(target_table.columns)):
            source_col_index = _source_index(col_index, source_col_count, repeat_after_first=False)
            _copy_cell_style(
                source_table.cell(source_row_index, source_col_index),
                target_table.cell(row_index, col_index),
            )


def _copy_dimensions(source_table: Any, target_table: Any, width: int, height: int) -> None:
    column_widths = [
        source_table.columns[_source_index(index, len(source_table.columns), repeat_after_first=False)].width
        for index in range(len(target_table.columns))
    ]
    row_heights = [
        source_table.rows[_source_index(index, len(source_table.rows), repeat_after_first=True)].height
        for index in range(len(target_table.rows))
    ]
    _assign_scaled_sizes(target_table.columns, column_widths, width)
    _assign_scaled_sizes(target_table.rows, row_heights, height)


def _assign_scaled_sizes(collection: Any, source_sizes: list[int], total_size: int) -> None:
    source_total = sum(source_sizes)
    if source_total <= 0:
        even_size = int(total_size / len(collection))
        for item in collection:
            if hasattr(item, "width"):
                item.width = even_size
            else:
                item.height = even_size
        return
    remaining = total_size
    last_index = len(source_sizes) - 1
    for index, item in enumerate(collection):
        if index == last_index:
            size = remaining
        else:
            size = int(total_size * source_sizes[index] / source_total)
            remaining -= size
        if hasattr(item, "width"):
            item.width = size
        else:
            item.height = size


def _source_index(index: int, source_count: int, repeat_after_first: bool) -> int:
    if source_count <= 1:
        return 0
    if index < source_count:
        return index
    if repeat_after_first:
        return 1 + ((index - 1) % (source_count - 1))
    return source_count - 1


def _copy_cell_style(source_cell: Any, target_cell: Any) -> None:
    source_tc = source_cell._tc
    target_tc = target_cell._tc
    _replace_xml_child(target_tc, target_tc.txBody, deepcopy(source_tc.txBody))
    _replace_xml_child(target_tc, target_tc.tcPr, deepcopy(source_tc.get_or_add_tcPr()))


def _replace_xml_child(parent: Any, old_child: Any, new_child: Any) -> None:
    index = parent.index(old_child)
    parent.remove(old_child)
    parent.insert(index, new_child)


def _apply_table_preserving_style(
    shape: Any,
    component: ComponentMapping,
    table_data: NormalizedTable,
) -> Any:
    table = shape.table
    row_count = table_data.row_count
    column_count = table_data.column_count
    if row_count > len(table.rows):
        reserved_rows = len(table.rows) - 1 if table_data.kind == "records" else len(table.rows)
        raise ReportGenerationError(
            ErrorCode.TABLE_OVERFLOW,
            f"组件 {component.location} 的数据行数 {table_data.data_row_count} 超过模板预留行数 {reserved_rows}",
            component,
        )
    if column_count > len(table.columns):
        raise ReportGenerationError(
            ErrorCode.TABLE_OVERFLOW,
            f"组件 {component.location} 的列数 {column_count} 超过模板预留列数 {len(table.columns)}",
            component,
        )

    _write_cells(table, table_data)
    for row_index in range(row_count, len(table.rows)):
        for col_index in range(len(table.columns)):
            _replace_cell_text_preserving_style(table.cell(row_index, col_index), "")
    for row_index in range(row_count):
        for col_index in range(column_count, len(table.columns)):
            _replace_cell_text_preserving_style(table.cell(row_index, col_index), "")

    return shape


def _write_cells(table: Any, table_data: NormalizedTable) -> None:
    for row_index in range(table_data.row_count):
        for col_index in range(table_data.column_count):
            cell = table.cell(row_index, col_index)
            style_cell = _fallback_style_cell(table, row_index, col_index)
            _replace_cell_text_preserving_style(cell, table_data.value_at(row_index, col_index), style_cell)


def _fallback_style_cell(table: Any, row_index: int, col_index: int) -> Any | None:
    for candidate_col in range(col_index - 1, -1, -1):
        candidate = table.cell(row_index, candidate_col)
        if _first_run(candidate) is not None:
            return candidate
    for candidate_row in range(row_index - 1, -1, -1):
        candidate = table.cell(candidate_row, col_index)
        if _first_run(candidate) is not None:
            return candidate
    return None


def _apply_table_placeholders(shape: Any, component: ComponentMapping, value: Any) -> Any:
    if not isinstance(value, Mapping):
        raise ReportGenerationError(
            ErrorCode.DATA_SOURCE_INVALID,
            f"组件 {component.location} 的占位符数据必须是对象",
            component,
        )

    table = shape.table
    for row_index in range(len(table.rows)):
        for col_index in range(len(table.columns)):
            cell = table.cell(row_index, col_index)
            template = cell.text
            if not _has_template_placeholder(template):
                continue
            rendered = _render_cell_template(component, template, value)
            style_cell = _fallback_style_cell(table, row_index, col_index)
            _replace_cell_text_preserving_style(
                cell,
                rendered,
                style_cell,
                prefer_fallback_for_unstyled=True,
            )
    return shape


def _has_template_placeholder(text: str) -> bool:
    return "{{" in text or "{%" in text


def _render_cell_template(component: ComponentMapping, template: str, context: Mapping[str, Any]) -> str:
    env = SandboxedEnvironment(undefined=StrictUndefined, autoescape=False)
    try:
        compiled = env.from_string(template)
    except TemplateSyntaxError as exc:
        raise ReportGenerationError(
            ErrorCode.DATA_SOURCE_INVALID,
            f"组件 {component.location} 的单元格占位符语法无效: {exc}",
            component,
        ) from exc
    try:
        return compiled.render(context)
    except UndefinedError as exc:
        raise ReportGenerationError(
            ErrorCode.DATA_SOURCE_NOT_FOUND,
            f"组件 {component.location} 的单元格占位符变量未找到: {exc}",
            component,
        ) from exc
    except TemplateError as exc:
        raise ReportGenerationError(
            ErrorCode.DATA_SOURCE_INVALID,
            f"组件 {component.location} 的单元格占位符渲染失败: {exc}",
            component,
        ) from exc


def _normalize_table(value: Any, component: ComponentMapping) -> NormalizedTable:
    if isinstance(value, dict) and "cells" in value:
        cells = _normalize_cells(value["cells"], component)
        return NormalizedTable(cells=cells, data_row_count=len(cells), kind="matrix")

    if isinstance(value, dict) and "columns" in value and "rows" in value:
        if not isinstance(value["columns"], list) or not isinstance(value["rows"], list):
            raise _invalid_table_data(component)
        for column in value["columns"]:
            if not isinstance(column, Mapping):
                raise _invalid_table_data(component)
        for row in value["rows"]:
            if not isinstance(row, Mapping):
                raise _invalid_table_data(component)
        columns = [
            {"key": str(column.get("key")), "label": str(column.get("label", column.get("key")))}
            for column in value["columns"]
        ]
        rows = [dict(row) for row in value["rows"]]
        return _records_table(columns, rows)

    if isinstance(value, list):
        for row in value:
            if not isinstance(row, Mapping):
                raise _invalid_table_data(component)
        rows = [dict(row) for row in value]
        order = component.config.get("order")
        if order:
            keys = [str(key) for key in order]
        else:
            keys = list(rows[0].keys()) if rows else []
        columns = [{"key": key, "label": key} for key in keys]
        return _records_table(columns, rows)

    raise ReportGenerationError(
        ErrorCode.DATA_SOURCE_INVALID,
        f"组件 {component.location} 的表格数据必须是对象数组、columns/rows 对象或 cells 矩阵对象",
        component,
    )


def _records_table(columns: list[dict[str, str]], rows: list[dict[str, Any]]) -> NormalizedTable:
    cells: list[list[Any]] = [[column["label"] for column in columns]]
    cells.extend([[row.get(column["key"], "") for column in columns] for row in rows])
    return NormalizedTable(cells=cells, data_row_count=len(rows), kind="records")


def _normalize_cells(value: Any, component: ComponentMapping) -> list[list[Any]]:
    if not isinstance(value, list):
        raise _invalid_table_data(component)
    cells: list[list[Any]] = []
    for row in value:
        if not isinstance(row, list):
            raise _invalid_table_data(component)
        cells.append(list(row))
    return cells


def _invalid_table_data(component: ComponentMapping) -> ReportGenerationError:
    return ReportGenerationError(
        ErrorCode.DATA_SOURCE_INVALID,
        f"组件 {component.location} 的表格数据必须是对象数组、columns/rows 对象或 cells 矩阵对象",
        component,
    )


def _table_font_size(component: ComponentMapping, row_count: int) -> int:
    minimum = int(component.config.get("min_font_size", 8))
    preferred = int(component.config.get("font_size", 12))
    if row_count <= 8:
        return preferred
    return max(minimum, preferred - (row_count - 8))


def _set_cell_text(cell: Any, value: Any, font_size: int, bold: bool = False) -> None:
    cell.text = "" if value is None else str(value)
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(font_size)
            run.font.bold = bold


def _apply_font_size(table: Any, font_size: int) -> None:
    for row in table.rows:
        for cell in row.cells:
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(font_size)


def _apply_column_value_styles(table: Any, component: ComponentMapping, table_data: NormalizedTable) -> None:
    styles = component.config.get("column_value_styles")
    if not isinstance(styles, Mapping) or not styles:
        return
    if table_data.kind != "records" or not table_data.cells:
        return

    headers = [str(value) for value in table_data.cells[0]]
    for column_name, value_styles in styles.items():
        if not isinstance(value_styles, Mapping):
            continue
        try:
            col_index = headers.index(str(column_name))
        except ValueError:
            continue
        for row_index in range(1, table_data.row_count):
            raw_value = table_data.value_at(row_index, col_index)
            style = value_styles.get(str(raw_value))
            if not isinstance(style, Mapping):
                continue
            _apply_cell_font_style(table.cell(row_index, col_index), style)


def _apply_cell_font_style(cell: Any, style: Mapping[str, Any]) -> None:
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            if "color" in style:
                run.font.color.rgb = _rgb(str(style["color"]))
            if "bold" in style:
                run.font.bold = bool(style["bold"])
            if "font_size" in style:
                run.font.size = Pt(int(style["font_size"]))


def _replace_cell_text_preserving_style(
    cell: Any,
    value: Any,
    style_cell: Any | None = None,
    *,
    prefer_fallback_for_unstyled: bool = False,
) -> None:
    text = "" if value is None else str(value)
    text_frame = cell.text_frame
    paragraph = text_frame.paragraphs[0]
    if paragraph.runs:
        if style_cell is not None and (
            _runs_are_empty(paragraph.runs)
            or (prefer_fallback_for_unstyled and not _run_has_style(paragraph.runs[0]))
        ):
            _copy_run_style(_first_run(style_cell), paragraph.runs[0])
        paragraph.runs[0].text = text
        for run in paragraph.runs[1:]:
            run.text = ""
    else:
        run = paragraph.add_run()
        _copy_run_style(_first_run(style_cell) if style_cell is not None else None, run)
        run.text = text
    for extra_paragraph in text_frame.paragraphs[1:]:
        for run in extra_paragraph.runs:
            run.text = ""


def _first_run(cell: Any) -> Any | None:
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            return run
    return None


def _runs_are_empty(runs: list[Any]) -> bool:
    return all(not run.text for run in runs)


def _run_has_style(run: Any) -> bool:
    run_properties = getattr(run._r, "rPr", None)
    if run_properties is None:
        return False
    return bool(run_properties.attrib) or len(run_properties) > 0


def _copy_run_style(source_run: Any | None, target_run: Any) -> None:
    if source_run is None:
        return
    source_r_pr = getattr(source_run._r, "rPr", None)
    if source_r_pr is None:
        return
    target_r = target_run._r
    target_r_pr = getattr(target_r, "rPr", None)
    if target_r_pr is not None:
        target_r.remove(target_r_pr)
    target_r.insert(0, deepcopy(source_r_pr))


def _rgb(value: str) -> RGBColor:
    normalized = value.lstrip("#")
    if len(normalized) != 6 or any(char not in "0123456789abcdefABCDEF" for char in normalized):
        raise ValueError(f"invalid RGB color {value!r}")
    return RGBColor(
        int(normalized[0:2], 16),
        int(normalized[2:4], 16),
        int(normalized[4:6], 16),
    )
