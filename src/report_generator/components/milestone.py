from __future__ import annotations

from collections.abc import Mapping
from typing import Any

from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_FILL
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from report_generator.errors import ErrorCode, ReportGenerationError
from report_generator.models import ComponentMapping
from report_generator.pptx.document import PptxDocument


DEFAULT_STATUS_STYLES: dict[str, dict[str, str]] = {
    "done": {"fill": "0066CC", "line": "0066CC", "text": "333333"},
    "active": {"fill": "FFFFFF", "line": "0066CC", "text": "333333"},
    "pending": {"fill": "FFFFFF", "line": "9AA6B2", "text": "555555"},
}

DEFAULT_HOLLOW_STATUSES = {"active", "pending"}


def apply_milestone(doc: PptxDocument, shape: Any, component: ComponentMapping, value: Any) -> None:
    try:
        items = _normalize_items(component, value)
        slide = shape.part.slide
        left = shape.left
        top = shape.top
        width = shape.width
        height = shape.height
        background_style = _extract_background_style(shape)
        preview_mode = str(component.config.get("preview_mode", "")).strip().lower()
        if preview_mode == "replace":
            _remove_preview_parts(doc, slide, component.location)
            if background_style:
                _promote_anchor_to_background(shape, component.location)
            else:
                doc.remove_shape(shape)
        elif preview_mode:
            raise ReportGenerationError(
                ErrorCode.DATA_SOURCE_INVALID,
                f"组件 {component.location} 当前只支持 preview_mode=replace",
                component,
            )
        elif bool(component.config.get("preserve_anchor", True)):
            _clear_shape_text(shape)
        else:
            doc.remove_shape(shape)
        if not items:
            return

        node_size = _emu(component.config.get("node_size"), Inches(0.13))
        line_y = top + int(height * float(component.config.get("line_y_ratio", 0.46)))
        date_height = _emu(component.config.get("date_height"), Inches(0.22))
        label_height = _emu(component.config.get("label_height"), Inches(0.24))
        font_size = int(component.config.get("font_size", 14))
        date_font_size = int(component.config.get("date_font_size", font_size))
        label_font_size = int(component.config.get("label_font_size", 14))
        date_font_name = str(component.config.get("date_font_name", component.config.get("font_name", "")))
        label_font_name = str(component.config.get("label_font_name", component.config.get("font_name", "Microsoft YaHei")))
        label_bold = bool(component.config.get("label_bold", True))
        node_inner_ratio = float(component.config.get("node_inner_ratio", 0.52))
        node_outline_width = float(component.config.get("node_outline_width", 0))
        default_text_axis_gap = node_size * 0.65 / Inches(1)
        text_axis_gap = _emu(component.config.get("text_axis_gap", default_text_axis_gap), int(node_size * 0.65))
        horizontal_padding = _emu(component.config.get("horizontal_padding"), 0)
        timeline_left = left + horizontal_padding
        timeline_width = max(1, width - horizontal_padding * 2)
        centers = _node_centers(timeline_left, timeline_width, len(items))
        text_slot_width = int(timeline_width / max(len(items), 1))
        date_width = min(width, _emu(component.config.get("date_width"), text_slot_width))

        if len(centers) > 1:
            line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, centers[0], line_y, centers[-1], line_y)
            line.name = f"{component.location}.line"
            line.line.color.rgb = _rgb(str(component.config.get("line_color", "222222")))
            line.line.width = Pt(float(component.config.get("line_width", 1.0)))

        for index, item in enumerate(items, start=1):
            center_x = centers[index - 1]
            styles = _item_style(component, item)
            _add_date(
                slide,
                f"{component.location}.item_{index}.date",
                str(item.get("date", "")),
                _centered_left(center_x, date_width),
                max(top, line_y - text_axis_gap - date_height),
                date_width,
                date_height,
                date_font_size,
                str(component.config.get("date_color", styles["text"])),
                date_font_name,
            )
            _add_node(
                slide,
                f"{component.location}.item_{index}.node",
                center_x - int(node_size / 2),
                line_y - int(node_size / 2),
                node_size,
                styles,
                hollow=_is_hollow_node(component, item, styles),
                inner_ratio=node_inner_ratio,
                outline_width=node_outline_width,
            )
            _add_label(
                slide,
                f"{component.location}.item_{index}.label",
                str(item.get("label", "")),
                _centered_left(center_x, text_slot_width),
                line_y + text_axis_gap,
                text_slot_width,
                label_height,
                label_font_size,
                styles["text"],
                label_bold,
                label_font_name,
            )
    except ReportGenerationError:
        raise
    except Exception as exc:
        raise ReportGenerationError(
            ErrorCode.DATA_SOURCE_INVALID,
            f"组件 {component.location} 的里程碑配置无效: {exc}",
            component,
        ) from exc


def _normalize_items(component: ComponentMapping, value: Any) -> list[dict[str, Any]]:
    if isinstance(value, Mapping) and "items" in value:
        value = value["items"]
    if not isinstance(value, list):
        raise ReportGenerationError(
            ErrorCode.DATA_SOURCE_INVALID,
            f"组件 {component.location} 的里程碑数据必须是数组或包含 items 的对象",
            component,
        )
    items: list[dict[str, Any]] = []
    for item in value:
        if not isinstance(item, Mapping):
            raise ReportGenerationError(
                ErrorCode.DATA_SOURCE_INVALID,
                f"组件 {component.location} 的里程碑条目必须是对象",
                component,
            )
        normalized = dict(item)
        normalized.setdefault("label", "")
        normalized.setdefault("date", "")
        normalized.setdefault("status", "pending")
        items.append(normalized)
    return items


def _clear_shape_text(shape: Any) -> None:
    if not getattr(shape, "has_text_frame", False):
        return
    text_frame = shape.text_frame
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.text = ""


def _extract_background_style(shape: Any) -> dict[str, Any]:
    style: dict[str, Any] = {}
    try:
        if shape.fill.type == MSO_FILL.SOLID:
            style["fill"] = True
    except (AttributeError, TypeError):
        pass
    try:
        if shape.line.fill.type == MSO_FILL.SOLID:
            style["line"] = True
    except (AttributeError, TypeError):
        pass
    return style


def _promote_anchor_to_background(shape: Any, component_location: str) -> None:
    shape.name = f"{component_location}.background"
    _clear_shape_text(shape)


def _remove_preview_parts(doc: PptxDocument, slide: Any, component_location: str) -> None:
    preview_prefix = f"{component_location}.preview."
    for preview_shape in list(slide.shapes):
        if str(getattr(preview_shape, "name", "")).startswith(preview_prefix):
            doc.remove_shape(preview_shape)


def _node_centers(left: int, width: int, count: int) -> list[int]:
    if count == 1:
        return [left + int(width / 2)]
    step = width / (count - 1)
    return [left + int(round(index * step)) for index in range(count)]


def _centered_left(center_x: int, width: int) -> int:
    return int(center_x - width / 2)


def _item_style(component: ComponentMapping, item: Mapping[str, Any]) -> dict[str, Any]:
    status = str(item.get("status", "pending"))
    style = dict(DEFAULT_STATUS_STYLES.get(status, DEFAULT_STATUS_STYLES["pending"]))
    configured_styles = component.config.get("status_styles", {})
    if isinstance(configured_styles, Mapping):
        configured = configured_styles.get(status)
        if isinstance(configured, Mapping):
            style.update({str(key): value for key, value in configured.items()})
    return style


def _is_hollow_node(component: ComponentMapping, item: Mapping[str, Any], styles: Mapping[str, Any]) -> bool:
    if "hollow" in styles:
        return _bool(styles["hollow"])
    configured = component.config.get("hollow_statuses")
    if isinstance(configured, list):
        return str(item.get("status", "pending")) in {str(value) for value in configured}
    return str(item.get("status", "pending")) in DEFAULT_HOLLOW_STATUSES


def _add_node(
    slide: Any,
    name: str,
    left: int,
    top: int,
    size: int,
    styles: Mapping[str, Any],
    *,
    hollow: bool,
    inner_ratio: float,
    outline_width: float,
) -> None:
    node_fill = str(styles["line"] if hollow else styles["fill"])
    node = _add_filled_circle(slide, name, left, top, size, node_fill, outline_width=outline_width)
    if not hollow:
        return

    inner_size = _same_parity(max(1, int(size * max(0.1, min(inner_ratio, 0.9)))), size)
    inner_left = int(left + (size - inner_size) / 2)
    inner_top = int(top + (size - inner_size) / 2)
    _add_filled_circle(slide, f"{name}.inner", inner_left, inner_top, inner_size, str(styles["fill"]), outline_width=outline_width)


def _add_filled_circle(
    slide: Any,
    name: str,
    left: int,
    top: int,
    size: int,
    color: str,
    *,
    outline_width: float,
) -> Any:
    node = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    node.name = name
    node.fill.solid()
    node.fill.fore_color.rgb = _rgb(color)
    if outline_width > 0:
        node.line.color.rgb = _rgb(color)
    node.line.width = Pt(outline_width)
    return node


def _add_date(
    slide: Any,
    name: str,
    text: str,
    left: int,
    top: int,
    width: int,
    height: int,
    font_size: int,
    color: str,
    font_name: str,
) -> None:
    _add_centered_text(slide, name, text, left, top, width, height, font_size, color, font_name=font_name, bold=True)


def _add_label(
    slide: Any,
    name: str,
    text: str,
    left: int,
    top: int,
    width: int,
    height: int,
    font_size: int,
    color: str,
    bold: bool,
    font_name: str,
) -> None:
    _add_centered_text(slide, name, text, left, top, width, height, font_size, color, font_name=font_name, bold=bold)


def _add_centered_text(
    slide: Any,
    name: str,
    text: str,
    left: int,
    top: int,
    width: int,
    height: int,
    font_size: int,
    color: str,
    *,
    font_name: str,
    bold: bool,
) -> None:
    textbox = slide.shapes.add_textbox(left, top, width, height)
    textbox.name = name
    text_frame = textbox.text_frame
    text_frame.clear()
    text_frame.margin_left = 0
    text_frame.margin_right = 0
    text_frame.margin_top = 0
    text_frame.margin_bottom = 0
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if font_name:
        run.font.name = font_name
    run.font.color.rgb = _rgb(color)


def _emu(value: Any, default: int) -> int:
    if value is None:
        return int(default)
    if isinstance(value, int | float):
        return int(Inches(value))
    if isinstance(value, str):
        try:
            return int(Inches(float(value)))
        except ValueError:
            return int(default)
    return int(default)


def _bool(value: Any) -> bool:
    if isinstance(value, bool):
        return value
    if isinstance(value, int | float):
        return bool(value)
    if isinstance(value, str):
        return value.strip().lower() in {"1", "true", "yes", "y", "on"}
    return bool(value)


def _same_parity(value: int, reference: int) -> int:
    if value % 2 == reference % 2:
        return value
    if value > 1:
        return value - 1
    return value + 1


def _rgb(value: str) -> RGBColor:
    normalized = str(value).lstrip("#")
    if len(normalized) != 6 or any(char not in "0123456789abcdefABCDEF" for char in normalized):
        raise ValueError(f"invalid RGB color {value!r}")
    return RGBColor(
        int(normalized[0:2], 16),
        int(normalized[2:4], 16),
        int(normalized[4:6], 16),
    )
