from __future__ import annotations

import math
from collections.abc import Mapping
from typing import Any

from pptx.dml.color import RGBColor
from pptx.util import Pt

from report_generator.errors import ErrorCode, ReportGenerationError
from report_generator.models import ComponentMapping

EMU_PER_INCH = 914400


def apply_text(shape: Any, component: ComponentMapping, value: Any) -> None:
    if not getattr(shape, "has_text_frame", False):
        raise ReportGenerationError(
            ErrorCode.TYPE_MISMATCH,
            f"组件 {component.location} 不是文本组件",
            component,
        )

    rich_text = _rich_text_runs(value)
    text = _rich_text_plain_text(rich_text) if rich_text is not None else ("" if value is None else str(value))
    min_font_size = int(component.config.get("min_font_size", 10))
    start_font_size = _existing_font_size(shape) or int(component.config.get("font_size", 18))
    if component.config.get("preserve_style"):
        if not _fits(shape, text, start_font_size):
            shape.height = _required_height(shape, text, start_font_size)
        if rich_text is not None:
            _apply_rich_text(shape, rich_text, preserve_style=True)
            return
        _apply_text_preserving_style(shape, text)
        return

    fitted_font_size = _fit_font_size(shape, text, start_font_size, min_font_size)
    if fitted_font_size is None:
        fitted_font_size = min_font_size
        shape.height = _required_height(shape, text, fitted_font_size)

    if rich_text is not None:
        _apply_rich_text(shape, rich_text, default_font_size=fitted_font_size)
        return

    text_frame = shape.text_frame
    text_frame.clear()
    lines = text.splitlines() or [""]
    for index, line in enumerate(lines):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        run = paragraph.add_run()
        run.text = line
        run.font.size = Pt(fitted_font_size)


def _rich_text_runs(value: Any) -> list[dict[str, Any]] | None:
    if not isinstance(value, Mapping):
        return None
    runs = value.get("rich_text")
    if runs is None:
        runs = value.get("runs")
    if runs is None:
        return None
    if not isinstance(runs, list):
        return None
    normalized: list[dict[str, Any]] = []
    for item in runs:
        if isinstance(item, Mapping):
            normalized.append(dict(item))
        else:
            normalized.append({"text": "" if item is None else str(item)})
    return normalized


def _rich_text_plain_text(runs: list[dict[str, Any]]) -> str:
    return "".join(str(run.get("text", "")) for run in runs)


def _apply_rich_text(
    shape: Any,
    runs: list[dict[str, Any]],
    *,
    preserve_style: bool = False,
    default_font_size: int | None = None,
) -> None:
    text_frame = shape.text_frame
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    for item in runs:
        pieces = str(item.get("text", "")).split("\n")
        for index, piece in enumerate(pieces):
            if index > 0:
                paragraph = text_frame.add_paragraph()
            run = paragraph.add_run()
            run.text = piece
            _apply_run_style(run, item, default_font_size=default_font_size, preserve_style=preserve_style)


def _apply_run_style(
    run: Any,
    style: Mapping[str, Any],
    *,
    default_font_size: int | None,
    preserve_style: bool,
) -> None:
    if "font_size" in style:
        run.font.size = Pt(int(style["font_size"]))
    elif default_font_size is not None and not preserve_style:
        run.font.size = Pt(default_font_size)
    if style.get("font_name"):
        run.font.name = str(style["font_name"])
    if "bold" in style:
        run.font.bold = _bool(style["bold"])
    if "italic" in style:
        run.font.italic = _bool(style["italic"])
    if "underline" in style:
        run.font.underline = _bool(style["underline"])
    if style.get("color"):
        run.font.color.rgb = _rgb(str(style["color"]))


def _apply_text_preserving_style(shape: Any, text: str) -> None:
    text_frame = shape.text_frame
    lines = text.splitlines() or [""]
    for paragraph_index, line in enumerate(lines):
        if paragraph_index < len(text_frame.paragraphs):
            paragraph = text_frame.paragraphs[paragraph_index]
        else:
            paragraph = text_frame.add_paragraph()
        if paragraph.runs:
            paragraph.runs[0].text = line
            for run in paragraph.runs[1:]:
                run.text = ""
        else:
            paragraph.add_run().text = line

    for paragraph in text_frame.paragraphs[len(lines) :]:
        for run in paragraph.runs:
            run.text = ""


def _existing_font_size(shape: Any) -> int | None:
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.size is not None:
                return int(round(run.font.size.pt))
    return None


def _fit_font_size(shape: Any, text: str, start: int, minimum: int) -> int | None:
    for font_size in range(start, minimum - 1, -1):
        if _fits(shape, text, font_size):
            return font_size
    return None


def _fits(shape: Any, text: str, font_size: int) -> bool:
    height_in = max(shape.height / EMU_PER_INCH, 0.1)
    needed_lines = _needed_line_count(shape, text, font_size)
    line_height_in = _line_height_in(font_size)
    capacity = max(1, math.floor(height_in / line_height_in))
    return needed_lines <= capacity


def _required_height(shape: Any, text: str, font_size: int) -> int:
    needed_lines = _needed_line_count(shape, text, font_size)
    line_height_in = _line_height_in(font_size)
    return max(shape.height, int(math.ceil(needed_lines * line_height_in * EMU_PER_INCH)))


def _needed_line_count(shape: Any, text: str, font_size: int) -> int:
    chars_per_line = _chars_per_line(shape, font_size)
    source_lines = text.splitlines() or [""]
    needed_lines = 0
    for line in source_lines:
        needed_lines += max(1, math.ceil(len(line) / chars_per_line))
    return needed_lines


def _chars_per_line(shape: Any, font_size: int) -> int:
    width_in = max(shape.width / EMU_PER_INCH, 0.1)
    return max(1, int((width_in * 72) / (font_size * 0.55)))


def _line_height_in(font_size: int) -> float:
    return (font_size * 1.25) / 72


def _bool(value: Any) -> bool:
    if isinstance(value, bool):
        return value
    if isinstance(value, int | float):
        return bool(value)
    if isinstance(value, str):
        return value.strip().lower() in {"1", "true", "yes", "y", "on"}
    return bool(value)


def _rgb(value: str) -> RGBColor:
    normalized = value.lstrip("#")
    if len(normalized) != 6 or any(char not in "0123456789abcdefABCDEF" for char in normalized):
        raise ValueError(f"invalid RGB color {value!r}")
    return RGBColor(
        int(normalized[0:2], 16),
        int(normalized[2:4], 16),
        int(normalized[4:6], 16),
    )
