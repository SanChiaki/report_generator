from __future__ import annotations

from collections.abc import Mapping
from typing import Any

from jinja2 import StrictUndefined
from jinja2.exceptions import TemplateError, TemplateSyntaxError, UndefinedError
from jinja2.sandbox import SandboxedEnvironment
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from report_generator.errors import ErrorCode, ReportGenerationError
from report_generator.models import ComponentMapping
from report_generator.pptx.document import PptxDocument


DEFAULT_STYLES: dict[str, dict[str, str]] = {
    "紧急": {"accent": "D64545", "fill": "FFFFFF", "line": "E6EAF0"},
    "重要": {"accent": "FF8A00", "fill": "FFFFFF", "line": "E6EAF0"},
    "一般": {"accent": "F5C400", "fill": "FFFFFF", "line": "E6EAF0"},
}


def apply_top_issues(doc: PptxDocument, shape: Any, component: ComponentMapping, value: Any) -> None:
    try:
        issues = _normalize_issues(component, value)
        slide = shape.part.slide
        left = shape.left
        top = shape.top
        width = shape.width
        card_height = _emu(component.config.get("card_height"), shape.height)
        card_gap = _emu(component.config.get("card_gap"), Inches(0.12))
        strip_width = _emu(component.config.get("strip_width"), Inches(0.08))
        severity_width = _emu(component.config.get("severity_width"), Inches(0.68))
        padding = _emu(component.config.get("padding"), Inches(0.12))
        preview_mode = str(component.config.get("preview_mode", "replace")).strip().lower()
        if preview_mode != "replace":
            raise ReportGenerationError(
                ErrorCode.DATA_SOURCE_INVALID,
                f"组件 {component.location} 当前只支持 preview_mode=replace",
                component,
            )

        _remove_preview_parts(doc, slide, component.location)
        doc.remove_shape(shape)

        for index, issue in enumerate(issues, start=1):
            item_top = top + (index - 1) * (card_height + card_gap)
            styles = _issue_style(component, issue)
            _add_card(slide, component.location, index, left, item_top, width, card_height, styles)
            _add_strip(slide, component.location, index, left, item_top, strip_width, card_height, styles)
            _add_text(
                slide,
                f"{component.location}.item_{index}.severity",
                str(issue.get("severity", "一般")),
                left + strip_width + padding,
                item_top + int(card_height * 0.16),
                severity_width,
                int(card_height * 0.36),
                font_size=int(component.config.get("severity_font_size", 18)),
                bold=True,
                color=styles["accent"],
            )

            text_left = left + strip_width + padding + severity_width
            text_width = width - strip_width - padding - severity_width - padding
            _add_text(
                slide,
                f"{component.location}.item_{index}.description",
                _render_issue_template(component, "description_template", _default_description(issue), issue),
                text_left,
                item_top + int(card_height * 0.12),
                text_width,
                int(card_height * 0.26),
                font_size=int(component.config.get("description_font_size", 13)),
                bold=True,
                color=str(component.config.get("description_color", "0B63CE")),
            )
            _add_text(
                slide,
                f"{component.location}.item_{index}.action",
                _render_issue_template(component, "action_template", "解决措施与进展：{{ action }}", issue),
                text_left,
                item_top + int(card_height * 0.42),
                text_width,
                int(card_height * 0.24),
                font_size=int(component.config.get("action_font_size", 13)),
                bold=True,
                color=str(component.config.get("action_color", "0B63CE")),
            )
            _add_text(
                slide,
                f"{component.location}.item_{index}.meta",
                _render_issue_template(component, "meta_template", _default_meta_template(), issue),
                text_left,
                item_top + int(card_height * 0.70),
                text_width,
                int(card_height * 0.22),
                font_size=int(component.config.get("meta_font_size", 12)),
                bold=False,
                color=str(component.config.get("meta_color", "222222")),
            )
    except ReportGenerationError:
        raise
    except Exception as exc:
        raise ReportGenerationError(
            ErrorCode.DATA_SOURCE_INVALID,
            f"组件 {component.location} 的 TOP 问题与风险配置无效: {exc}",
            component,
        ) from exc


def _normalize_issues(component: ComponentMapping, value: Any) -> list[dict[str, Any]]:
    if isinstance(value, Mapping) and "items" in value:
        value = value["items"]
    if not isinstance(value, list):
        raise ReportGenerationError(
            ErrorCode.DATA_SOURCE_INVALID,
            f"组件 {component.location} 的 TOP 问题与风险数据必须是数组或包含 items 的对象",
            component,
        )

    issues: list[dict[str, Any]] = []
    for item in value:
        if not isinstance(item, Mapping):
            raise ReportGenerationError(
                ErrorCode.DATA_SOURCE_INVALID,
                f"组件 {component.location} 的 TOP 问题与风险条目必须是对象",
                component,
            )
        normalized = dict(item)
        normalized.setdefault("severity", "一般")
        normalized.setdefault("created_at", "")
        normalized.setdefault("description", normalized.get("title", ""))
        normalized.setdefault("action", "")
        normalized.setdefault("owner", "")
        normalized.setdefault("status", "")
        normalized.setdefault("due_date", "")
        issues.append(normalized)
    return issues


def _remove_preview_parts(doc: PptxDocument, slide: Any, component_location: str) -> None:
    preview_prefix = f"{component_location}.preview."
    for preview_shape in list(slide.shapes):
        if str(getattr(preview_shape, "name", "")).startswith(preview_prefix):
            doc.remove_shape(preview_shape)


def _issue_style(component: ComponentMapping, issue: Mapping[str, Any]) -> dict[str, str]:
    severity = str(issue.get("severity", "一般"))
    style = dict(DEFAULT_STYLES.get(severity, DEFAULT_STYLES["一般"]))
    configured_styles = component.config.get("styles", {})
    if isinstance(configured_styles, Mapping):
        configured = configured_styles.get(severity)
        if isinstance(configured, Mapping):
            style.update({str(key): str(value) for key, value in configured.items()})
    return style


def _add_card(
    slide: Any,
    base_name: str,
    index: int,
    left: int,
    top: int,
    width: int,
    height: int,
    styles: Mapping[str, str],
) -> None:
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    card.name = f"{base_name}.item_{index}.card"
    card.fill.solid()
    card.fill.fore_color.rgb = _rgb(styles["fill"])
    card.line.color.rgb = _rgb(styles["line"])


def _add_strip(
    slide: Any,
    base_name: str,
    index: int,
    left: int,
    top: int,
    width: int,
    height: int,
    styles: Mapping[str, str],
) -> None:
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    strip.name = f"{base_name}.item_{index}.strip"
    strip.fill.solid()
    strip.fill.fore_color.rgb = _rgb(styles["accent"])
    strip.line.color.rgb = _rgb(styles["accent"])


def _add_text(
    slide: Any,
    name: str,
    text: str,
    left: int,
    top: int,
    width: int,
    height: int,
    *,
    font_size: int,
    bold: bool,
    color: str,
) -> None:
    textbox = slide.shapes.add_textbox(left, top, width, height)
    textbox.name = name
    text_frame = textbox.text_frame
    text_frame.clear()
    text_frame.margin_left = 0
    text_frame.margin_right = 0
    text_frame.margin_top = 0
    text_frame.margin_bottom = 0
    paragraph = text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)


def _render_issue_template(
    component: ComponentMapping,
    config_key: str,
    default_template: str,
    issue: Mapping[str, Any],
) -> str:
    template = str(component.config.get(config_key, default_template))
    env = SandboxedEnvironment(undefined=StrictUndefined, autoescape=False)
    try:
        compiled = env.from_string(template)
    except TemplateSyntaxError as exc:
        raise ReportGenerationError(
            ErrorCode.DATA_SOURCE_INVALID,
            f"组件 {component.location} 的 {config_key} 模板语法无效: {exc}",
            component,
        ) from exc
    try:
        return compiled.render(dict(issue))
    except UndefinedError as exc:
        raise ReportGenerationError(
            ErrorCode.DATA_SOURCE_NOT_FOUND,
            f"组件 {component.location} 的 {config_key} 模板变量未找到: {exc}",
            component,
        ) from exc
    except TemplateError as exc:
        raise ReportGenerationError(
            ErrorCode.DATA_SOURCE_INVALID,
            f"组件 {component.location} 的 {config_key} 模板渲染失败: {exc}",
            component,
        ) from exc


def _default_description(issue: Mapping[str, Any]) -> str:
    if issue.get("created_at"):
        return "问题描述：创建时间{{ created_at }}，当前进展{{ description }}"
    return "问题描述：{{ description }}"


def _default_meta_template() -> str:
    return "责任人：{{ owner }}   |  状态：{{ status }}  | 计划解决时间：{{ due_date }}"


def _emu(value: Any, default: int) -> int:
    if value is None:
        return int(default)
    if isinstance(value, int | float):
        return int(Inches(value))
    if isinstance(value, str):
        try:
            return int(Inches(float(value)))
        except ValueError:
            return int(default)
    return int(default)


def _rgb(value: str) -> RGBColor:
    normalized = str(value).lstrip("#")
    if len(normalized) != 6 or any(char not in "0123456789abcdefABCDEF" for char in normalized):
        raise ValueError(f"invalid RGB color {value!r}")
    return RGBColor(
        int(normalized[0:2], 16),
        int(normalized[2:4], 16),
        int(normalized[4:6], 16),
    )
