from __future__ import annotations

from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


NAVY = "22324D"
TEAL = "0F766E"
GREEN = "16A34A"
AMBER = "F59E0B"
RED = "DC2626"
LIGHT_BG = "F6F8FB"
LIGHT_BLUE = "E8F1FF"
LIGHT_GREEN = "EAF7EF"
LIGHT_AMBER = "FFF4D8"
TEXT = "111827"
MUTED = "6B7280"
WHITE = "FFFFFF"


def main() -> None:
    base_dir = Path(__file__).resolve().parent
    logo_path = base_dir / "company_logo.png"
    _ensure_logo(logo_path)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    _add_cover_slide(prs)
    _add_risk_slide(prs)
    _add_delivery_slide(prs)

    output = base_dir / "sample_template.pptx"
    prs.save(output)
    print(output)


def _add_cover_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide)

    title = _add_text(slide, "text.report_title", "项目月报标题", 0.55, 0.45, 7.2, 0.55, 30, NAVY, True)
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    _add_text(slide, "text.report_period", "报告周期｜负责人", 0.58, 1.05, 5.6, 0.35, 12, MUTED)
    _add_text(slide, "text.executive_summary", "本月项目摘要，将由 JSON 写入。", 0.58, 1.55, 7.0, 1.05, 18, TEXT)

    logo = _add_text(slide, "image.company_logo", "Logo", 10.65, 0.45, 2.25, 0.75, 18, MUTED, True)
    _style_box(logo, "FFFFFF", "D1D5DB")

    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.15), Inches(1.55), Inches(2.0), Inches(0.5))
    badge.name = "shape.project_status_badge"
    badge.text = "整体状态"
    _style_box(badge, GREEN, GREEN)
    _style_text_frame(badge, 15, WHITE, True, PP_ALIGN.CENTER)

    _add_metric_card(slide, "text.kpi_progress", "完成率", "0%", 0.58, 3.0, LIGHT_GREEN)
    _add_metric_card(slide, "text.kpi_budget", "预算使用率", "0%", 3.28, 3.0, LIGHT_BLUE)
    _add_metric_card(slide, "text.kpi_risks", "风险数量", "0", 5.98, 3.0, LIGHT_AMBER)

    table = slide.shapes.add_table(2, 4, Inches(0.58), Inches(4.8), Inches(11.9), Inches(1.55))
    table.name = "table.key_metrics"
    _seed_table(table.table, ["指标", "当前值", "目标值", "说明"], ["示例指标", "0", "0", "示例说明"])


def _add_risk_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide)

    _add_text(slide, "text.risk_title", "风险与里程碑", 0.55, 0.42, 6.0, 0.55, 27, NAVY, True)
    _add_text(slide, "text.risk_summary", "本页展示 TOP 风险和关键里程碑。", 0.58, 1.02, 7.2, 0.65, 16, TEXT)

    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.75), Inches(0.52), Inches(1.75), Inches(0.48))
    badge.name = "shape.risk_status_badge"
    badge.text = "风险状态"
    _style_box(badge, AMBER, AMBER)
    _style_text_frame(badge, 13, WHITE, True, PP_ALIGN.CENTER)

    top_risks = slide.shapes.add_table(3, 4, Inches(0.58), Inches(1.75), Inches(11.9), Inches(2.35))
    top_risks.name = "table.top_risks"
    _seed_table(
        top_risks.table,
        ["风险类型", "风险等级", "责任人", "风险描述"],
        ["示例类型", "中", "示例负责人", "示例描述"],
    )

    milestones = slide.shapes.add_table(3, 4, Inches(0.58), Inches(4.55), Inches(11.9), Inches(1.65))
    milestones.name = "table.milestones"
    _seed_table(
        milestones.table,
        ["阶段", "计划日期", "状态", "完成说明"],
        ["示例阶段", "2026-06-01", "进行中", "示例说明"],
    )


def _add_delivery_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide)

    _add_text(slide, "text.delivery_title", "交付趋势与下月动作", 0.55, 0.42, 6.0, 0.55, 27, NAVY, True)
    _add_text(slide, "text.next_step_summary", "下月重点动作摘要。", 0.6, 6.38, 8.0, 0.45, 14, TEXT)

    progress_chart = _add_chart(slide, "chart.progress_trend", "完成率", 0.58, 1.35, 5.75, 3.0)
    progress_chart.chart.has_legend = False
    budget_chart = _add_chart(slide, "chart.budget_usage", "预算使用率", 6.75, 1.35, 5.75, 3.0)
    budget_chart.chart.has_legend = False

    _add_metric_card(slide, "text.quality_score", "质量得分", "0", 0.58, 4.75, LIGHT_GREEN)
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.3), Inches(5.02), Inches(1.7), Inches(0.5))
    badge.name = "shape.quality_badge"
    badge.text = "质量状态"
    _style_box(badge, GREEN, GREEN)
    _style_text_frame(badge, 13, WHITE, True, PP_ALIGN.CENTER)

    actions = slide.shapes.add_table(2, 3, Inches(5.35), Inches(4.75), Inches(7.15), Inches(1.35))
    actions.name = "table.next_actions"
    _seed_table(actions.table, ["动作", "负责人", "截止时间"], ["示例动作", "示例负责人", "2026-06-30"])


def _add_background(slide) -> None:
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    _style_box(bg, LIGHT_BG, LIGHT_BG)


def _add_metric_card(slide, name: str, label: str, value: str, x: float, y: float, fill: str) -> None:
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.25), Inches(1.12))
    _style_box(card, fill, fill)
    _add_text(slide, f"label.{name}", label, x + 0.18, y + 0.18, 1.8, 0.25, 11, MUTED)
    value_box = _add_text(slide, name, value, x + 0.18, y + 0.48, 1.75, 0.45, 25, NAVY, True)
    value_box.text_frame.margin_left = 0


def _add_chart(slide, name: str, series_name: str, x: float, y: float, w: float, h: float):
    chart_data = CategoryChartData()
    chart_data.categories = ["一月", "二月", "三月"]
    chart_data.add_series(series_name, (10, 20, 30))
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
        chart_data,
    )
    chart_shape.name = name
    chart_shape.chart.chart_title.has_text_frame = True
    chart_shape.chart.chart_title.text_frame.text = series_name
    return chart_shape


def _add_text(
    slide,
    name: str,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: int,
    color: str,
    bold: bool = False,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    shape.name = name
    shape.text_frame.clear()
    paragraph = shape.text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)
    return shape


def _seed_table(table, headers: list[str], first_row: list[str]) -> None:
    for index, header in enumerate(headers):
        _set_cell(table.cell(0, index), header, 10, True, WHITE)
        table.cell(0, index).fill.solid()
        table.cell(0, index).fill.fore_color.rgb = _rgb(TEAL)
    for index, value in enumerate(first_row):
        _set_cell(table.cell(1, index), value, 9, False, TEXT)
    for index in range(len(headers)):
        table.columns[index].width = int(table.columns[index].width)


def _set_cell(cell, value: str, size: int, bold: bool, color: str) -> None:
    cell.text = value
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = _rgb(color)


def _style_box(shape, fill: str, line: str) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill)
    shape.line.color.rgb = _rgb(line)


def _style_text_frame(shape, size: int, color: str, bold: bool, align) -> None:
    for paragraph in shape.text_frame.paragraphs:
        paragraph.alignment = align
        for run in paragraph.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = _rgb(color)


def _ensure_logo(path: Path) -> None:
    image = Image.new("RGB", (900, 300), "#FFFFFF")
    draw = ImageDraw.Draw(image)
    draw.rounded_rectangle((18, 18, 882, 282), radius=42, fill="#22324D")
    draw.rounded_rectangle((48, 58, 210, 220), radius=28, fill="#0F766E")
    try:
        title_font = ImageFont.truetype("/System/Library/Fonts/Supplemental/Arial Bold.ttf", 92)
        small_font = ImageFont.truetype("/System/Library/Fonts/Supplemental/Arial.ttf", 42)
    except OSError:
        title_font = ImageFont.load_default()
        small_font = ImageFont.load_default()
    draw.text((84, 82), "RG", fill="#FFFFFF", font=title_font)
    draw.text((255, 72), "REPORT", fill="#FFFFFF", font=title_font)
    draw.text((260, 180), "GENERATOR", fill="#D1FAE5", font=small_font)
    image.save(path)


def _rgb(hex_color: str) -> RGBColor:
    normalized = hex_color.lstrip("#")
    return RGBColor(int(normalized[:2], 16), int(normalized[2:4], 16), int(normalized[4:], 16))


if __name__ == "__main__":
    main()
