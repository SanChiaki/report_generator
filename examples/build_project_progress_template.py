from __future__ import annotations

import json
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from report_generator.generator import generate_report


BASE_DIR = Path(__file__).resolve().parent
TEMPLATE_PATH = BASE_DIR / "project_progress_template.pptx"
REPORT_PATH = BASE_DIR / "project_progress_report.pptx"
MAPPING_PATH = BASE_DIR / "project_progress_mapping.json"
PAYLOAD_PATH = BASE_DIR / "project_progress_payload.json"

INK = "1F2937"
SLATE = "475569"
BLUE = "2563EB"
BLUE_DARK = "1E3A8A"
CYAN = "0891B2"
GREEN = "16A34A"
AMBER = "F59E0B"
RED = "DC2626"
BG = "F8FAFC"
PANEL = "FFFFFF"
LINE = "CBD5E1"
TABLE_ALT = "EEF4FF"


def main() -> None:
    prs = build_template()
    prs.save(TEMPLATE_PATH)
    print(TEMPLATE_PATH)


def generate_example_report() -> None:
    result = generate_report(
        TEMPLATE_PATH.read_bytes(),
        json.loads(MAPPING_PATH.read_text(encoding="utf-8")),
        json.loads(PAYLOAD_PATH.read_text(encoding="utf-8")),
    )
    REPORT_PATH.write_bytes(result)
    print(REPORT_PATH)


def build_template() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    _cover(prs.slides.add_slide(blank))
    _project_info(prs.slides.add_slide(blank))
    _dashboard(prs.slides.add_slide(blank))
    _overall_progress(prs.slides.add_slide(blank))
    _delivery_plan(prs.slides.add_slide(blank))
    _task_details(prs.slides.add_slide(blank))
    _issues_and_risks(prs.slides.add_slide(blank))
    _current_progress(prs.slides.add_slide(blank))
    return prs


def _cover(slide) -> None:
    _background(slide, dark=True)
    _decorative_bar(slide, 0.0, 0.0, 13.333, 0.18, CYAN)
    _decorative_bar(slide, 0.0, 7.24, 13.333, 0.26, BLUE)
    _text(slide, "text.cover_title", "xxx项目进展报告", 0.78, 2.35, 8.5, 0.75, 34, "FFFFFF", True)
    _text(
        slide,
        "text.cover_period",
        "报告周期：2025-12-14至2025-12-25",
        0.82,
        3.35,
        6.2,
        0.38,
        16,
        "DCEBFF",
    )
    _text(slide, "text.cover_report_date", "报告日期：2025-12-25", 0.82, 3.9, 4.0, 0.36, 13, "BBD7FF")
    _text(slide, "static.cover_label", "PROJECT PROGRESS", 9.5, 0.72, 2.5, 0.32, 11, "93C5FD", True)
    _circle(slide, 10.1, 2.25, 1.15, CYAN)
    _circle(slide, 11.0, 2.85, 1.85, BLUE)
    _circle(slide, 9.25, 3.6, 0.75, GREEN)


def _project_info(slide) -> None:
    _section_slide(slide, "项目信息", "基础项目信息由 JSON 写入，模板字段和样式保持固定。")
    shape = _table(slide, "table.project_info", 8, 2, 0.82, 1.52, 11.7, 4.95, ["字段", "内容"])
    labels = ["项目名称", "项目编号", "客户名称", "项目经理", "计划开始时间", "计划结束时间", "报告日期"]
    for row, label in enumerate(labels, start=1):
        _cell(shape.table.cell(row, 0), label, 13, INK, bold=True, fill="EAF2FF")
        _cell(shape.table.cell(row, 1), "待填充", 13, INK, fill=PANEL)


def _dashboard(slide) -> None:
    _section_slide(slide, "项目仪表盘", "关键状态汇总。")
    shape = _table(slide, "table.project_dashboard", 7, 2, 0.82, 1.52, 11.7, 4.9, ["指标", "状态"])
    labels = ["项目状态", "项目进度", "进度状态", "技术状态", "资源状态", "经营状态"]
    for row, label in enumerate(labels, start=1):
        _cell(shape.table.cell(row, 0), label, 14, INK, bold=True, fill="EAF2FF")
        _cell(shape.table.cell(row, 1), "待填充", 14, INK, fill=PANEL)


def _overall_progress(slide) -> None:
    _section_slide(slide, "项目整体进展", "本页保留模板中的摘要文本样式和进展表样式。")
    _text(
        slide,
        "text.overall_progress_summary",
        "项目整体进展摘要。",
        0.82,
        1.45,
        11.7,
        0.72,
        16,
        INK,
        fill=PANEL,
        line=LINE,
        margin=0.15,
    )
    shape = _table(slide, "table.overall_progress", 5, 4, 0.82, 2.45, 11.7, 3.55, ["进展项", "完成情况", "状态", "说明"])
    for row in range(1, 5):
        _cell(shape.table.cell(row, 0), "待填充", 12, INK, fill=PANEL if row % 2 else TABLE_ALT)
        _cell(shape.table.cell(row, 1), "待填充", 12, INK, fill=PANEL if row % 2 else TABLE_ALT)
        _cell(shape.table.cell(row, 2), "待填充", 12, INK, fill=PANEL if row % 2 else TABLE_ALT)
        _cell(shape.table.cell(row, 3), "待填充", 12, INK, fill=PANEL if row % 2 else TABLE_ALT)


def _delivery_plan(slide) -> None:
    _section_slide(slide, "项目交付计划", "交付计划表预留 6 条计划行。")
    shape = _table(
        slide,
        "table.delivery_plan",
        7,
        5,
        0.62,
        1.45,
        12.1,
        5.25,
        ["阶段", "计划开始", "计划结束", "交付物", "状态"],
    )
    for row in range(1, 7):
        fill = PANEL if row % 2 else TABLE_ALT
        for col in range(5):
            _cell(shape.table.cell(row, col), "待填充", 11, INK, fill=fill)


def _task_details(slide) -> None:
    _section_slide(slide, "任务详情", "阶段任务责任关系。")
    shape = _table(
        slide,
        "table.task_details",
        9,
        4,
        0.62,
        1.35,
        12.1,
        5.65,
        ["阶段名称", "任务名称", "任务状态", "责任人"],
    )
    for row in range(1, 9):
        fill = PANEL if row % 2 else TABLE_ALT
        for col in range(4):
            _cell(shape.table.cell(row, col), "待填充", 11, INK, fill=fill)


def _issues_and_risks(slide) -> None:
    _section_slide(slide, "TOP 问题及风险", "高优先级问题与风险跟踪。")
    shape = _table(
        slide,
        "table.top_issues_risks",
        7,
        5,
        0.62,
        1.35,
        12.1,
        5.4,
        ["类型", "等级", "描述", "责任人", "应对措施"],
    )
    for row in range(1, 7):
        fill = PANEL if row % 2 else "FFF7ED"
        for col in range(5):
            _cell(shape.table.cell(row, col), "待填充", 10, INK, fill=fill)


def _current_progress(slide) -> None:
    _section_slide(slide, "当前进展", "本期完成、下期计划和当前阻塞。")
    _text(
        slide,
        "text.current_progress_summary",
        "当前进展摘要。",
        0.82,
        1.35,
        11.7,
        0.78,
        16,
        INK,
        fill=PANEL,
        line=LINE,
        margin=0.15,
    )
    shape = _table(slide, "table.current_progress", 5, 3, 0.82, 2.45, 11.7, 3.55, ["类别", "事项", "状态"])
    for row in range(1, 5):
        fill = PANEL if row % 2 else TABLE_ALT
        for col in range(3):
            _cell(shape.table.cell(row, col), "待填充", 12, INK, fill=fill)


def _section_slide(slide, title: str, subtitle: str) -> None:
    _background(slide)
    _decorative_bar(slide, 0.0, 0.0, 13.333, 0.16, BLUE)
    _text(slide, f"static.title.{title}", title, 0.62, 0.42, 5.8, 0.48, 24, INK, True)
    _text(slide, f"static.subtitle.{title}", subtitle, 0.64, 0.96, 8.4, 0.28, 10, SLATE)
    _footer(slide)


def _background(slide, dark: bool = False) -> None:
    fill = BLUE_DARK if dark else BG
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.name = "static.background"
    _style_shape(bg, fill, fill)


def _footer(slide) -> None:
    _text(slide, "static.footer", "Project Progress Report", 0.64, 7.08, 2.7, 0.18, 8, "94A3B8")
    _decorative_bar(slide, 10.7, 7.16, 1.9, 0.04, CYAN)


def _table(slide, name: str, rows: int, cols: int, x: float, y: float, w: float, h: float, headers: list[str]):
    shape = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.name = name
    for col in range(cols):
        shape.table.columns[col].width = int(Inches(w) / cols)
    for row in range(rows):
        shape.table.rows[row].height = int(Inches(h) / rows)
    for col, header in enumerate(headers):
        _cell(shape.table.cell(0, col), header, 12, "FFFFFF", bold=True, fill=BLUE)
    return shape


def _text(
    slide,
    name: str,
    content: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: int,
    color: str,
    bold: bool = False,
    fill: str | None = None,
    line: str | None = None,
    margin: float = 0.0,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    shape.name = name
    if fill:
        _style_shape(shape, fill, line or fill)
    shape.text_frame.clear()
    shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    shape.text_frame.margin_left = Inches(margin)
    shape.text_frame.margin_right = Inches(margin)
    paragraph = shape.text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = content
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)
    return shape


def _cell(cell, text: str, size: int, color: str, bold: bool = False, fill: str = PANEL) -> None:
    cell.fill.solid()
    cell.fill.fore_color.rgb = _rgb(fill)
    cell.margin_left = Inches(0.08)
    cell.margin_right = Inches(0.08)
    cell.margin_top = Inches(0.04)
    cell.margin_bottom = Inches(0.04)
    cell.text = text
    for paragraph in cell.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.LEFT
        for run in paragraph.runs:
            run.font.name = "Arial"
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = _rgb(color)


def _decorative_bar(slide, x: float, y: float, w: float, h: float, color: str) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.name = "static.decorative_bar"
    _style_shape(shape, color, color)


def _circle(slide, x: float, y: float, size: float, color: str) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size))
    shape.name = "static.circle"
    _style_shape(shape, color, color)


def _style_shape(shape, fill: str, line: str) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill)
    shape.line.color.rgb = _rgb(line)


def _rgb(hex_color: str) -> RGBColor:
    value = hex_color.lstrip("#")
    return RGBColor(int(value[:2], 16), int(value[2:4], 16), int(value[4:6], 16))


if __name__ == "__main__":
    main()
